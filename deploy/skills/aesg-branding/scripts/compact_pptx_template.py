#!/usr/bin/env python3
"""Create the compact AESG runtime PPTX from the approved 17-slide specimen deck."""

from __future__ import annotations

import argparse
import json
import os
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation


EXPECTED_SIZE_EMU = (9906000, 6858000)


def package_is_clean(path: Path) -> None:
    with zipfile.ZipFile(path) as archive:
        corrupt = archive.testzip()
    if corrupt:
        raise ValueError(f"corrupt ZIP member: {corrupt}")


def hierarchy(presentation: Presentation) -> dict:
    return {
        "slides": len(presentation.slides),
        "masters": len(presentation.slide_masters),
        "layouts": sum(len(master.slide_layouts) for master in presentation.slide_masters),
        "sizeEmu": [presentation.slide_width, presentation.slide_height],
        "slideLayouts": [slide.slide_layout.name for slide in presentation.slides],
    }


def compact(source: Path, output: Path) -> dict:
    package_is_clean(source)
    presentation = Presentation(source)
    before = hierarchy(presentation)

    used_layout_parts = {slide.slide_layout.part for slide in presentation.slides}
    used_master_parts = {slide.slide_layout.slide_master.part for slide in presentation.slides}

    for master in list(presentation.slide_masters):
        layout_ids = master._element.sldLayoutIdLst
        for layout_id in list(layout_ids):
            relationship = master.part.rels[layout_id.rId]
            if relationship.target_part in used_layout_parts:
                continue
            master.part.drop_rel(layout_id.rId)
            layout_ids.remove(layout_id)

    master_ids = presentation._element.sldMasterIdLst
    for master_id in list(master_ids):
        relationship = presentation.part.rels[master_id.rId]
        if relationship.target_part in used_master_parts:
            continue
        presentation.part.drop_rel(master_id.rId)
        master_ids.remove(master_id)

    presentation.core_properties.subject = "AESG Compact General Template"
    output.parent.mkdir(parents=True, exist_ok=True)
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{output.stem}-",
        suffix=output.suffix,
        dir=output.parent,
    )
    os.close(descriptor)
    temporary = Path(temporary_name)
    try:
        presentation.save(temporary)
        package_is_clean(temporary)
        check = Presentation(temporary)
        after = hierarchy(check)
        if after["slides"] != 17:
            raise ValueError(f"expected 17 specimen slides, found {after['slides']}")
        if after["masters"] != 1 or after["layouts"] != 17:
            raise ValueError(
                "compact runtime must contain exactly one master and seventeen layouts"
            )
        if tuple(after["sizeEmu"]) != EXPECTED_SIZE_EMU:
            raise ValueError(f"unexpected presentation size: {after['sizeEmu']}")
        temporary.replace(output)
    finally:
        temporary.unlink(missing_ok=True)

    return {
        "source": str(source),
        "output": str(output),
        "bytes": output.stat().st_size,
        "before": before,
        "after": after,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    source = args.source.resolve()
    output = args.output.resolve()
    if source == output:
        raise ValueError("source and output must be different files")
    if source.suffix.casefold() != ".pptx" or output.suffix.casefold() != ".pptx":
        raise ValueError("source and output must both use the .pptx extension")
    if not source.is_file():
        raise FileNotFoundError(f"source presentation not found: {source}")
    print(json.dumps(compact(source, output), indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
