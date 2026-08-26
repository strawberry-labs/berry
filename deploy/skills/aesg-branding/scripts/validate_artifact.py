#!/usr/bin/env python3
"""Fast structural checks for AESG PDF and OpenXML artifacts."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree


PLACEHOLDER_RE = re.compile(
    r"(?:\b(?:lorem ipsum|insert title|sample heading|sample sub-heading|"
    r"replace text|name another|name surname|client name|click to add|"
    r"click to edit|section heading|xxxx-xxx)\b|x{8,}|"
    r"\[[^\]]+\|\s*(?:word|excel|powerpoint)\])",
    re.IGNORECASE,
)
ERROR_TOKENS = ("#REF!", "#DIV/0!", "#VALUE!", "#NAME?", "#N/A", "#NUM!", "#NULL!")
WORD_TEXT_PART_RE = re.compile(
    r"word/(?:document|header\d+|footer\d+|footnotes|endnotes|comments)\.xml$"
)
SLIDE_PART_RE = re.compile(r"ppt/slides/slide(\d+)\.xml$")
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

PPT_BODY_SHAPE_IDS = {
    499, 500, 503, 504, 507, 508, 515, 518, 522, 524, 528, 533, 537, 541,
    546, 547, 548, 550, 551, 552, 553, 554, 557,
}
PPT_TITLE_SHAPE_IDS = {498, 502, 506, 509, 517, 521, 527, 531, 556}
PPT_SECTION_SHAPE_IDS = {497, 501, 505, 516, 520, 526, 530}
PPT_DIVIDER_SHAPE_IDS = {532}


def run(command: list[str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, capture_output=True, text=True, check=False)


def pptx_rgb(font) -> str | None:
    try:
        value = font.color.rgb
    except ValueError:
        return None
    return str(value).upper() if value is not None else None


def package_text(path: Path) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(path) as archive:
        corrupt = archive.testzip()
        if corrupt:
            raise ValueError(f"corrupt ZIP member: {corrupt}")
        for name in archive.namelist():
            if name.endswith((".xml", ".rels")):
                chunks.append(archive.read(name).decode("utf-8", errors="ignore"))
    return "\n".join(chunks)


def paragraph_texts(xml: bytes, namespace: str) -> list[str]:
    root = ElementTree.fromstring(xml)
    paragraph_tag = f"{{{namespace}}}p"
    text_tag = f"{{{namespace}}}t"
    paragraphs: list[str] = []
    for paragraph in root.iter(paragraph_tag):
        if any(
            descendant is not paragraph and descendant.tag == paragraph_tag
            for descendant in paragraph.iter()
        ):
            continue
        paragraphs.append(
            "".join(node.text or "" for node in paragraph.iter(text_tag))
        )
    return paragraphs


def docx_story_text(path: Path) -> str:
    chunks: list[str] = []
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if WORD_TEXT_PART_RE.fullmatch(name):
                chunks.extend(
                    paragraph_texts(archive.read(name), WORD_NS)
                )
    return "\n".join(chunks)


def inspect_pptx_slides(
    path: Path, slide_order: dict[str, int]
) -> tuple[str, list[str], list[str]]:
    text: list[str] = []
    empty_placeholders: list[str] = []
    picture_placeholders: list[str] = []
    namespaces = {"p": PRESENTATION_NS, "a": DRAWING_NS}
    with zipfile.ZipFile(path) as archive:
        slide_parts = sorted(
            (
                (slide_order.get(name, int(match.group(1))), name)
                for name in archive.namelist()
                if (match := SLIDE_PART_RE.fullmatch(name))
            ),
            key=lambda item: item[0],
        )
        for slide_number, name in slide_parts:
            slide_xml = archive.read(name)
            root = ElementTree.fromstring(slide_xml)
            slide_text = paragraph_texts(slide_xml, DRAWING_NS)
            text.extend(slide_text)
            for shape in root.findall(".//p:sp", namespaces):
                properties = shape.find("./p:nvSpPr/p:cNvPr", namespaces)
                shape_name = properties.get("name", "") if properties is not None else ""
                placeholder = shape.find("./p:nvSpPr/p:nvPr/p:ph", namespaces)
                shape_text = "".join(
                    node.text or "" for node in shape.findall(".//a:t", namespaces)
                ).strip()
                if "picture placeholder" in shape_name.casefold():
                    picture_placeholders.append(f"slide {slide_number}, {shape_name}")
                if placeholder is None or shape_text:
                    continue
                placeholder_type = placeholder.get("type", "obj")
                if placeholder_type not in {"dt", "ftr"}:
                    empty_placeholders.append(
                        f"slide {slide_number}, {shape_name or placeholder_type}"
                    )
    return "\n".join(text), empty_placeholders, picture_placeholders


def validate_core_author(path: Path, errors: list[str], evidence: dict) -> None:
    with zipfile.ZipFile(path) as archive:
        try:
            core = archive.read("docProps/core.xml").decode("utf-8", errors="ignore")
        except KeyError:
            errors.append("Office package has no core properties")
            return
    creators = re.findall(r"<(?:dc:creator|cp:lastModifiedBy)>(.*?)</", core)
    evidence["coreAuthors"] = creators
    if creators and any(value.strip().casefold() != "aesg" for value in creators):
        errors.append("Office core properties contain a non-AESG author")


def validate_pdf(path: Path, errors: list[str], evidence: dict) -> None:
    if path.read_bytes()[:5] != b"%PDF-":
        errors.append("missing PDF signature")
    if shutil.which("qpdf"):
        result = run(["qpdf", "--check", str(path)])
        evidence["qpdf"] = result.stdout.strip() or result.stderr.strip()
        if result.returncode:
            errors.append("qpdf validation failed")
    try:
        from pypdf import PdfReader

        document = PdfReader(path)
        extracted = "\n".join(page.extract_text() or "" for page in document.pages)
        placeholders = sorted(set(match.group(0) for match in PLACEHOLDER_RE.finditer(extracted)))
        if placeholders:
            errors.append(f"placeholder text remains: {', '.join(placeholders)}")
        evidence["pypdf"] = {"pages": len(document.pages), "textChecked": True}
    except Exception as exc:
        errors.append(f"PDF package validation failed: {exc}")
    if shutil.which("pdfinfo"):
        result = run(["pdfinfo", str(path)])
        evidence["pdfinfo"] = result.stdout.strip() or result.stderr.strip()
        if result.returncode:
            errors.append("pdfinfo validation failed")
    else:
        errors.append("missing validator: pdfinfo")
    if shutil.which("pdffonts"):
        fonts = run(["pdffonts", str(path)])
        evidence["pdffonts"] = fonts.stdout.strip()
        font_text = fonts.stdout.casefold()
        evidence["fontContract"] = {
            "verdana": "verdana" in font_text,
            "fallbacks": [
                name
                for name in (
                    "linuxlibertine",
                    "liberationserif",
                    "liberationsans",
                    "dejavuserif",
                    "dejavusans",
                    "noto serif",
                    "noto sans",
                )
                if name in font_text
            ],
        }
        if "verdana" not in font_text:
            errors.append("Verdana is not embedded in AESG PDF")
        if evidence["fontContract"]["fallbacks"]:
            errors.append(
                "fallback fonts embedded in AESG PDF: "
                + ", ".join(evidence["fontContract"]["fallbacks"])
            )
    if shutil.which("file"):
        mime = run(["file", "--brief", "--mime-type", str(path)]).stdout.strip()
        evidence["mime"] = mime
        if mime != "application/pdf":
            errors.append(f"unexpected MIME type: {mime}")


def validate_docx(path: Path, errors: list[str], evidence: dict) -> None:
    from docx import Document
    from docx.oxml.ns import qn

    document = Document(path)
    text = docx_story_text(path)
    placeholders = sorted(set(match.group(0) for match in PLACEHOLDER_RE.finditer(text)))
    if placeholders:
        errors.append(f"placeholder text remains: {', '.join(placeholders)}")
    broken_fields = [
        token
        for token in ("error! bookmark not defined", "includepicture", "click here to enter text")
        if token in text.casefold()
    ]
    if broken_fields:
        errors.append("broken or specimen fields remain: " + ", ".join(broken_fields))
    if not document.sections:
        errors.append("DOCX has no sections")
    xml = package_text(path).casefold()
    if "verdana" not in xml:
        errors.append("Verdana is not declared in DOCX package")
    with zipfile.ZipFile(path) as archive:
        styles_xml = archive.read("word/styles.xml").decode("utf-8", errors="ignore")
        theme_xml = archive.read("word/theme/theme1.xml").decode("utf-8", errors="ignore")
        theme_font_attrs = sorted(
            set(
                re.findall(
                    r"w:(?:asciiTheme|hAnsiTheme|eastAsiaTheme|csTheme|cstheme)=\"([^\"]+)\"",
                    styles_xml,
                )
            )
        )
        theme_typefaces = sorted(
            set(
                re.findall(
                    r"<a:(?:latin|ea|cs)\b[^>]*typeface=\"([^\"]*)\"",
                    theme_xml,
                )
            )
        )
        evidence["docxTheme"] = {
            "styleThemeFonts": theme_font_attrs,
            "typefaces": theme_typefaces,
        }
        if theme_font_attrs:
            errors.append(
                "DOCX styles still contain theme font mappings: "
                + ", ".join(theme_font_attrs)
            )
        if any(value and value.casefold() != "verdana" for value in theme_typefaces):
            errors.append("DOCX theme contains a non-Verdana typeface")
        footer_highlights = [
            name
            for name in archive.namelist()
            if re.fullmatch(r"word/footer\d+\.xml", name)
            and b"<w:highlight" in archive.read(name)
        ]
        if footer_highlights:
            errors.append("DOCX footer highlighting remains: " + ", ".join(footer_highlights))
    style_map = {style.name.casefold(): style for style in document.styles if style.name}

    def check_style(name: str, size_pt: float, color: str) -> None:
        style = style_map.get(name.casefold())
        if style is None:
            errors.append(f"missing required AESG style: {name}")
            return
        if (style.font.name or "").casefold() != "verdana":
            errors.append(f"{name} is not explicitly Verdana")
        actual_size = style.font.size.pt if style.font.size is not None else None
        if actual_size is None or abs(actual_size - size_pt) > 0.05:
            errors.append(f"{name} size is {actual_size}, expected {size_pt} pt")
        actual_color = style.font.color.rgb
        if actual_color is None or str(actual_color).upper() != color:
            errors.append(f"{name} colour is {actual_color}, expected {color}")

    is_report = "aesg main heading" in style_map
    if is_report:
        check_style("Normal", 10, "343741")
        check_style("AESG Body Text", 10, "343741")
        check_style("AESG Main Heading", 22, "008C95")
        check_style("AESG Sub H1", 16, "008C95")
        check_style("AESG H2", 14, "008C95")
        check_style("AESG Bullet", 10, "343741")
        check_style("AESG Numbering", 10, "343741")
        check_style("AESG Table Caption", 9, "04999A")
        check_style("AESG Figure Captions", 9, "008C95")
        content_sections = document.sections[1:] if len(document.sections) > 1 else document.sections
        top_margins = [round(section.top_margin.twips) for section in content_sections]
        evidence["reportTopMarginsDxa"] = top_margins
        expected_margins = [
            720 if section.page_width > section.page_height else 1701
            for section in content_sections
        ]
        if top_margins != expected_margins:
            errors.append(
                f"report content top margins are {top_margins}; expected {expected_margins} DXA"
            )
        heading_styles = {"aesg main heading", "aesg sub h1", "aesg h2"}
        manual_heading_re = re.compile(
            r"^\s*(?:\d+(?:\.\d+)*(?:[.)])?\s+|appendix\s+[a-z]\b)",
            re.IGNORECASE,
        )

        def numbering_disabled(paragraph) -> bool:
            properties = paragraph._p.find(qn("w:pPr"))
            numbering = properties.find(qn("w:numPr")) if properties is not None else None
            num_id = numbering.find(qn("w:numId")) if numbering is not None else None
            return num_id is not None and num_id.get(qn("w:val")) == "0"

        duplicate_numbering = [
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
            and paragraph.style is not None
            and paragraph.style.name.casefold() in heading_styles
            and manual_heading_re.match(paragraph.text)
            and not numbering_disabled(paragraph)
        ]
        if duplicate_numbering:
            errors.append(
                "source-numbered AESG headings still inherit automatic numbering: "
                + ", ".join(duplicate_numbering[:5])
            )

        list_intro_labels = {
            "recommendations:",
            "examples:",
            "suggested priorities:",
            "the locking strategy should include:",
        }
        list_style_names = {"aesg bullet", "list bullet", "aesg numbering", "list number"}
        suspicious_lists: list[str] = []
        paragraphs = document.paragraphs
        for index, paragraph in enumerate(paragraphs):
            label = paragraph.text.strip().casefold()
            if label not in list_intro_labels:
                continue
            candidates = []
            for following in paragraphs[index + 1 :]:
                text_value = following.text.strip()
                if not text_value:
                    continue
                style_name = following.style.name.casefold() if following.style else ""
                if style_name in heading_styles or text_value.casefold() in list_intro_labels:
                    break
                if len(text_value) <= 180:
                    candidates.append(following)
                if len(candidates) >= 5 or len(text_value) > 180:
                    break
            if len(candidates) >= 3 and not any(
                item.style is not None and item.style.name.casefold() in list_style_names
                for item in candidates
            ):
                suspicious_lists.append(paragraph.text.strip())
        if suspicious_lists:
            errors.append(
                "list-shaped report blocks are flattened as body paragraphs after: "
                + ", ".join(suspicious_lists[:5])
            )

        empty_compliance_rows = 0
        rag_values: list[str] = []
        for table in document.tables:
            if not table.rows:
                continue
            header = [cell.text.strip().casefold() for cell in table.rows[0].cells]
            is_compliance = "requirement" in header and ("gap" in header or "rag" in header)
            if not is_compliance:
                continue
            for row in table.rows[1:]:
                values = [cell.text.strip() for cell in row.cells]
                if all(not value for value in values):
                    empty_compliance_rows += 1
                if values and values[-1].casefold() in {"green", "amber", "red"}:
                    rag_values.append(values[-1])
        if empty_compliance_rows:
            errors.append(
                f"compliance table contains {empty_compliance_rows} empty scaffold rows"
            )
        if any(value not in {"Green", "Amber", "Red"} for value in rag_values):
            errors.append("compliance-table RAG values are not consistently title-cased")
        evidence["reportStructure"] = {
            "sourceNumberedHeadings": len(duplicate_numbering),
            "flattenedListBlocks": len(suspicious_lists),
            "emptyComplianceRows": empty_compliance_rows,
            "ragValues": sorted(set(rag_values)),
        }
    else:
        check_style("Heading 1", 12, "008C95")
        check_style("Heading 2", 9, "53565A")
        check_style("Normal", 9, "53565A")
    validate_core_author(path, errors, evidence)
    evidence.update(
        {
            "sections": len(document.sections),
            "paragraphs": len(document.paragraphs),
            "tables": len(document.tables),
        }
    )


def validate_xlsx(path: Path, errors: list[str], evidence: dict) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False)
    if not workbook.sheetnames:
        errors.append("workbook has no worksheets")
    hidden = [ws.title for ws in workbook.worksheets if ws.sheet_state != "visible"]
    if hidden:
        errors.append(f"unexpected hidden worksheets: {', '.join(hidden)}")
    if any(name.casefold() == "joiners" for name in workbook.sheetnames):
        errors.append("forbidden Joiners worksheet found")
    formula_errors: list[str] = []
    placeholders: list[str] = []
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                value = cell.value
                if isinstance(value, str):
                    if any(token in value for token in ERROR_TOKENS):
                        formula_errors.append(f"{worksheet.title}!{cell.coordinate}")
                    if PLACEHOLDER_RE.search(value):
                        placeholders.append(f"{worksheet.title}!{cell.coordinate}")
    if formula_errors:
        errors.append(f"formula error tokens at: {', '.join(formula_errors[:12])}")
    if placeholders:
        errors.append(f"placeholder text at: {', '.join(placeholders[:12])}")
    xml = package_text(path).casefold()
    if "joiners" in xml:
        errors.append("forbidden Joiners identifier found in package")
    if "verdana" not in xml:
        errors.append("Verdana is not declared in XLSX package")
    validate_core_author(path, errors, evidence)
    evidence["worksheets"] = workbook.sheetnames
    evidence["dimensions"] = {
        ws.title: {"rows": ws.max_row, "columns": ws.max_column} for ws in workbook.worksheets
    }


def validate_pptx(path: Path, errors: list[str], evidence: dict) -> None:
    from pptx import Presentation

    presentation = Presentation(path)
    slide_order = {
        str(slide.part.partname).lstrip("/"): index
        for index, slide in enumerate(presentation.slides, start=1)
    }
    slide_text, empty_placeholders, picture_placeholders = inspect_pptx_slides(
        path, slide_order
    )
    placeholders = sorted(set(match.group(0) for match in PLACEHOLDER_RE.finditer(slide_text)))
    if placeholders:
        errors.append(f"placeholder text remains: {', '.join(placeholders[:12])}")
    if empty_placeholders:
        errors.append(
            f"empty structural placeholders remain: {', '.join(empty_placeholders[:12])}"
        )
    if picture_placeholders:
        errors.append(
            f"unresolved picture placeholders remain: {', '.join(picture_placeholders[:12])}"
        )
    if not presentation.slides:
        errors.append("presentation has no slides")
    xml = package_text(path).casefold()
    if "verdana" not in xml:
        errors.append("Verdana is not declared in PPTX package")
    typography_errors: list[str] = []
    role_contract = {
        "body": (PPT_BODY_SHAPE_IDS, 9.0, "343741"),
        "title": (PPT_TITLE_SHAPE_IDS, 21.0, "343741"),
        "section": (PPT_SECTION_SHAPE_IDS, 8.5, "008C95"),
        "divider": (PPT_DIVIDER_SHAPE_IDS, 22.0, "FFFFFF"),
    }
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            role = next(
                (
                    role_name
                    for role_name, (shape_ids, _size, _color) in role_contract.items()
                    if shape.shape_id in shape_ids
                ),
                None,
            )
            if role is None:
                continue
            _shape_ids, expected_size, expected_color = role_contract[role]
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text:
                        continue
                    font_name = (run.font.name or "").casefold()
                    actual_size = run.font.size.pt if run.font.size is not None else None
                    actual_color = pptx_rgb(run.font)
                    if font_name != "verdana":
                        typography_errors.append(f"slide {slide_number} shape {shape.shape_id} {role}: font {run.font.name}")
                    if actual_size is None or abs(actual_size - expected_size) > 0.05:
                        typography_errors.append(f"slide {slide_number} shape {shape.shape_id} {role}: size {actual_size}")
                    if actual_color is None or str(actual_color).upper() != expected_color:
                        typography_errors.append(f"slide {slide_number} shape {shape.shape_id} {role}: colour {actual_color}")
    if typography_errors:
        errors.append("PPTX typography contract failed: " + "; ".join(typography_errors[:12]))
    layout_count = sum(len(master.slide_layouts) for master in presentation.slide_masters)
    expected_size = [9906000, 6858000]
    actual_size = [presentation.slide_width, presentation.slide_height]
    if presentation.core_properties.subject == "AESG Compact General Template":
        if len(presentation.slide_masters) != 1 or layout_count != 17:
            errors.append("AESG compact runtime must contain one master and seventeen layouts")
        if actual_size != expected_size:
            errors.append(f"unexpected AESG compact template size: {actual_size}")
    validate_core_author(path, errors, evidence)
    evidence.update(
        {
            "slides": len(presentation.slides),
            "layouts": layout_count,
            "masters": len(presentation.slide_masters),
            "sizeEmu": actual_size,
            "emptyPlaceholders": empty_placeholders,
            "picturePlaceholders": picture_placeholders,
            "typographyContract": "passed" if not typography_errors else typography_errors[:12],
        }
    )


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("artifact", type=Path)
    args = parser.parse_args()
    path = args.artifact.resolve()
    errors: list[str] = []
    evidence: dict = {"path": str(path), "bytes": path.stat().st_size if path.is_file() else 0}

    if not path.is_file() or path.stat().st_size == 0:
        errors.append("artifact is missing or empty")
    else:
        try:
            if path.suffix.casefold() == ".pdf":
                validate_pdf(path, errors, evidence)
            elif path.suffix.casefold() == ".docx":
                validate_docx(path, errors, evidence)
            elif path.suffix.casefold() == ".xlsx":
                validate_xlsx(path, errors, evidence)
            elif path.suffix.casefold() == ".pptx":
                validate_pptx(path, errors, evidence)
            else:
                errors.append(f"unsupported extension: {path.suffix}")
        except Exception as exc:
            errors.append(f"validation exception: {exc}")

    report = {"ok": not errors, "errors": errors, "evidence": evidence}
    print(json.dumps(report, indent=2))
    return 0 if report["ok"] else 1


if __name__ == "__main__":
    sys.exit(main())
