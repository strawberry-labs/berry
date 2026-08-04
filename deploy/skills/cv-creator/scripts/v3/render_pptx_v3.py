#!/usr/bin/env python3
"""Render V3 AESG PPTX CVs from the preprocessed V3 CV JSON.

This renderer is PPTX-only and does not touch the DOCX V3 workflow. It keeps
the existing PPTX templates as the design source, then paginates variable
content into card-sized chunks so text is not truncated.
"""

from __future__ import annotations

import argparse
import base64
import copy
import io
import json
import math
import re
import subprocess
import sys
import tempfile
from collections import Counter
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from cv_template_common import (  # noqa: E402
    artifact_output_path,
    delete_pptx_slide,
    fit_profile_name_lines,
    slugify,
)
from text_normalization import repair_text_artifacts  # noqa: E402


ROOT = Path(__file__).resolve().parents[2]
DATA_PATH = ROOT / "output" / "extracted" / "preprocessed_cv_data.json"
TEMPLATE_DIR = ROOT / "assets" / "templates" / "v3"
PORTRAIT_TEMPLATE = TEMPLATE_DIR / "aesg_cv_portrait_v3.pptx"
LANDSCAPE_TEMPLATE = TEMPLATE_DIR / "aesg_cv_landscape_v3.pptx"
WHITE_LOGO_ASSET = TEMPLATE_DIR / "aesg_logo_white.png"
PORTRAIT_OUTPUT = ROOT / "output" / "generated" / "v3" / "pptx_portrait"
LANDSCAPE_OUTPUT = ROOT / "output" / "generated" / "v3" / "pptx_landscape"

AESG_TEAL = "008C95"
TEXT_DARK = "303030"
MUTED_GREY = "9A9DA9"
WHITE = "FFFFFF"

FONT_BODY = "Ubuntu"
FONT_TITLE = "Roboto Slab"

EMU_PER_INCH = 914400
LAYOUT_LABELS = {"KEY EXPERTISE", "WORK EXPERIENCE", "QUALIFICATIONS", "SELECTED PROJECTS"}
EXPERIENCE_SEPARATOR = "\ue000"


@dataclass
class TextStyle:
    heading_size: float
    body_size: float
    project_title_size: float
    project_body_size: float
    overview_size: float
    role_size: float
    name_size: float
    experience_body_start: int
    body_color: str = TEXT_DARK
    body_font: str = FONT_BODY
    title_font: str = FONT_TITLE


@dataclass
class ProjectChunk:
    project: dict[str, Any]
    body: str
    part: int
    parts: int


@dataclass(frozen=True)
class SplitContinuationLayout:
    section_left: int
    divider_x: int
    project_left: int
    top: int
    bottom: int
    right: int
    project_columns: int = 1
    column_gap: int = 150000
    draw_divider: bool = True
    project_only_full_width: bool = False
    masonry_projects: bool = False
    section_density_factor: float = 98.0


PORTRAIT_STYLE = TextStyle(
    # Portrait Version 1 uses Verdana. The original placeholders inherit this
    # from the slide master, so every generated run sets it explicitly rather
    # than relying on the viewer's font substitution.
    heading_size=10.0,
    body_size=9.0,
    project_title_size=9.0,
    project_body_size=9.0,
    overview_size=9.0,
    role_size=9.0,
    name_size=28.0,
    experience_body_start=int(1.18 * EMU_PER_INCH),
    body_font="Verdana",
    title_font="Verdana",
)

LANDSCAPE_STYLE = TextStyle(
    # The new landscape source deck uses the same Verdana family as the
    # portrait deck.  Set it on every generated run; relying on the deck's
    # theme substitutes Arial/Ubuntu differently across PowerPoint viewers.
    # The new source deck uses a 9pt Verdana text system.  The previous
    # compact values were inherited from the retired landscape artwork and
    # caused viewers to show 7.5pt/7.7pt body text.
    heading_size=10.0,
    body_size=9.0,
    project_title_size=9.0,
    project_body_size=9.0,
    overview_size=9.0,
    role_size=9.0,
    name_size=28.0,
    experience_body_start=int(0.98 * EMU_PER_INCH),
    body_font="Verdana",
    title_font="Verdana",
)


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def clean_text(value: str) -> str:
    value = repair_text_artifacts(str(value or ""))
    value = value.replace("\u00a0", " ")
    value = re.sub(r"\s+", " ", value).strip()
    return value


def project_client_display(project: dict[str, Any]) -> str:
    return clean_text(
        project.get("client_line") or project.get("client") or "Confidential"
    )


def line_items(items: Iterable[Any]) -> list[str]:
    result: list[str] = []
    for item in items or []:
        text = item.get("text") if isinstance(item, dict) else str(item)
        text = clean_text(text)
        if text:
            result.append(text)
    return result


def experience_line_items(items: Iterable[Any]) -> list[str]:
    """Keep the date prefix separate from the role/company body for layout."""
    result: list[str] = []
    for item in items or []:
        if isinstance(item, dict):
            prefix = clean_text(item.get("prefix") or "")
            body = clean_text(item.get("body") or item.get("text") or "")
        else:
            prefix = ""
            body = clean_text(str(item))
        if not body and not prefix:
            continue
        result.append(f"{prefix}{EXPERIENCE_SEPARATOR}{body}" if prefix else body)
    return result


def load_cvs(path: Path) -> list[dict[str, Any]]:
    with path.open(encoding="utf-8") as handle:
        payload = json.load(handle)
    return payload.get("cvs", [payload])


def clear_shape(shape: Any) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    # Remove any paragraph borders inherited from the placeholder so that
    # section headings (e.g., MEMBERSHIPS) don't keep a top horizontal line.
    A_PPR = "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    A_PBDR = "{http://schemas.openxmlformats.org/drawingml/2006/main}pBdr"
    for paragraph in tf.paragraphs:
        p_pr = paragraph._p.find(A_PPR)
        if p_pr is not None:
            p_bdr = p_pr.find(A_PBDR)
            if p_bdr is not None:
                p_pr.remove(p_bdr)
    # Explicitly remove inherited placeholder styling. Google Slides otherwise
    # renders some empty/filled placeholders as white cards with a drop shadow.
    try:
        shape.fill.background()
        shape.line.fill.background()
    except Exception:
        pass
    sp_pr = getattr(shape.element, "spPr", None)
    if sp_pr is not None:
        for child in list(sp_pr):
            if child.tag.endswith("effectLst") or child.tag.endswith("effectDag"):
                sp_pr.remove(child)


def add_paragraph(
    shape: Any,
    text: str = "",
    *,
    size: float,
    bold: bool = False,
    color: str = TEXT_DARK,
    font: str = FONT_BODY,
    bullet: bool = False,
    space_after: float = 0,
    level: int = 0,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    tf = shape.text_frame
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
        paragraph = tf.paragraphs[0]
    else:
        paragraph = tf.add_paragraph()
    # `TextFrame.clear()` retains the first placeholder paragraph's `pPr`.
    # The portrait source uses bulleted placeholders, which otherwise gives a
    # generated heading an inherited bullet even though the heading itself is
    # not a list item. All renderer bullets are explicit teal runs, so remove
    # any inherited bullet definition before applying our paragraph styling.
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        tag_name = child.tag.rsplit("}", 1)[-1]
        if tag_name in {
            "buChar", "buAutoNum", "buBlip", "buNone", "buFont", "buClr",
        }:
            p_pr.remove(child)
        if tag_name == "pBdr":
            # Remove any inherited paragraph borders (e.g., a top horizontal
            # line above the first reused placeholder paragraph) so headings
            # render uniformly across left-sidebar sections.
            p_pr.remove(child)
    paragraph.alignment = align
    paragraph.level = level
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(space_after)
    paragraph.line_spacing = 1.0
    paragraph.font.size = Pt(size)
    paragraph.font.name = font
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    paragraph.text = ""
    if bullet:
        try:
            paragraph._p.get_or_add_pPr().set("marL", "171450")
            paragraph._p.get_or_add_pPr().set("indent", "-91440")
        except Exception:
            pass
        text = f"\u2022 {text}"
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return paragraph


def fill_plain(
    shape: Any,
    text: str,
    *,
    size: float,
    bold: bool = False,
    color: str = TEXT_DARK,
    font: str = FONT_BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    clear_shape(shape)
    add_paragraph(shape, clean_text(text), size=size, bold=bold, color=color, font=font, align=align)


def fill_name(
    shape: Any,
    first_name: str,
    last_name: str,
    *,
    size: float,
    color: str = WHITE,
    font: str = FONT_TITLE,
) -> None:
    """Fill a profile name using no more than two explicit lines.

    The source template uses a *single paragraph* with a soft line break
    between first and last name.  Two paragraphs introduce paragraph spacing
    and push the candidate title into the second line. Preserve that structure,
    but rebalance and fit long names before disabling automatic word wrapping.
    """
    # Materialise the placeholder's inherited layout geometry on the slide.
    # Without a local xfrm, LibreOffice and padded-canvas validators can leave
    # the text at its master position while moving the rest of the slide,
    # which makes a two-line name appear above the page boundary.
    set_shape_geometry(shape, shape.left, shape.top, shape.width, shape.height)
    clear_shape(shape)
    # The landscape master supplies these values through its layout.  Once a
    # placeholder is populated directly, python-pptx otherwise falls back to
    # PowerPoint's default text-box insets, which makes the two-line name sit
    # lower and too close to the candidate title.  Preserve the master’s
    # zero-inset, 80% line-rhythm explicitly.
    text_frame = shape.text_frame
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    font_path = (
        "/System/Library/Fonts/Supplemental/Verdana Bold.ttf"
        if font.lower() == "verdana"
        else "/System/Library/Fonts/Supplemental/Verdana Bold.ttf"
    )

    def measure_name(text: str, size_pt: float) -> float:
        try:
            from PIL import ImageFont

            measured_font = ImageFont.truetype(
                font_path,
                max(1, round(size_pt * 96 / 72)),
            )
            return float(measured_font.getlength(text))
        except (OSError, AttributeError):
            return len(text) * size_pt * 0.72

    line_one, line_two, fitted_size = fit_profile_name_lines(
        first_name,
        last_name,
        max_width=max(1.0, shape.width / EMU_PER_INCH * 96),
        measure=measure_name,
        base_size=size,
    )
    paragraph = text_frame.paragraphs[0]
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.rsplit("}", 1)[-1] in {
            "buChar", "buAutoNum", "buBlip", "buNone", "buFont", "buClr",
        }:
            p_pr.remove(child)
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    # PowerPoint's source placeholder uses an 80% soft-break line rhythm.
    paragraph.line_spacing = 0.80
    paragraph.text = ""
    parts = [part for part in (line_one, line_two) if part]
    for index, part in enumerate(parts):
        if index:
            # python-pptx does not serialise a literal vertical-tab correctly
            # in a run (it becomes `_x000B_` in PowerPoint).  Emit its native
            # DrawingML break element instead.
            paragraph.add_line_break()
        name_run = paragraph.add_run()
        name_run.text = part
        name_run.font.size = Pt(fitted_size)
        name_run.font.name = font
        name_run.font.bold = True
        name_run.font.color.rgb = rgb(color)


def append_project_text(
    shape: Any,
    chunk: ProjectChunk,
    style: TextStyle,
    *,
    gap_after: float = 0,
) -> None:
    """Append one editable project to an existing flowing text frame."""
    project = chunk.project
    title = clean_text(project.get("name") or "Project")
    if chunk.parts > 1 and chunk.part > 1:
        title = f"{title} (continued {chunk.part}/{chunk.parts})"
    client = project_client_display(project)
    duration = clean_text(project.get("duration_line") or project.get("duration") or "")

    paragraph = add_paragraph(
        shape,
        "",
        size=style.project_title_size,
        bold=True,
        color=AESG_TEAL,
        font=style.title_font,
        space_after=2,
    )
    title_run = paragraph.add_run()
    title_run.text = title
    title_run.font.name = style.title_font
    title_run.font.size = Pt(style.project_title_size)
    title_run.font.bold = True
    title_run.font.color.rgb = rgb(AESG_TEAL)
    last_paragraph = paragraph
    if client and chunk.part == 1:
        last_paragraph = add_paragraph(
            shape,
            client,
            size=style.project_body_size,
            bold=False,
            color=TEXT_DARK,
            font=style.body_font,
            space_after=1,
        )
    if duration and chunk.part == 1:
        last_paragraph = add_paragraph(
            shape,
            duration,
            size=style.project_body_size,
            bold=False,
            color=TEXT_DARK,
            font=style.body_font,
            space_after=2,
        )
    if chunk.body:
        last_paragraph = add_paragraph(
            shape,
            chunk.body,
            size=style.project_body_size,
            color=TEXT_DARK,
            font=style.body_font,
            align=PP_ALIGN.JUSTIFY,
        )
    # In a flowing project textbox, paragraph spacing is the sole source of
    # inter-project whitespace. Exactly one body-text line is therefore
    # applied to the final paragraph of every project except the last.
    last_paragraph.space_after = Pt(gap_after)


def add_project_text(shape: Any, chunk: ProjectChunk, style: TextStyle) -> None:
    clear_shape(shape)
    append_project_text(shape, chunk, style)


def add_projects_flowing_text(shape: Any, chunks: list[ProjectChunk], style: TextStyle) -> None:
    """Render multiple projects into one native editable PowerPoint textbox."""
    clear_shape(shape)
    for index, chunk in enumerate(chunks):
        append_project_text(
            shape,
            chunk,
            style,
            gap_after=style.project_body_size if index < len(chunks) - 1 else 0,
        )


def add_section_text(shape: Any, section_chunks: list[tuple[str, list[str], bool]], style: TextStyle) -> None:
    clear_shape(shape)
    for index, (heading, items, continued) in enumerate(section_chunks):
        # A continued section flows from the preceding slide without repeating
        # an artificial "CONT." label. New sections still receive one heading.
        if not continued:
            add_paragraph(
                shape,
                heading.upper(),
                size=style.heading_size,
                bold=True,
                color=AESG_TEAL,
                font=style.title_font,
                space_after=2,
            )
        for item in items:
            if heading.casefold() == "experience":
                add_experience_paragraph(shape, item, style)
            else:
                add_sidebar_item_paragraph(shape, item, style)
        if index < len(section_chunks) - 1:
            add_paragraph(shape, "", size=style.body_size, color=AESG_TEAL, font=style.body_font, space_after=3)


def add_items_text(shape: Any, items: list[str], style: TextStyle) -> None:
    clear_shape(shape)
    for item in items:
        if EXPERIENCE_SEPARATOR in item:
            add_experience_paragraph(shape, item, style)
        else:
            add_sidebar_item_paragraph(shape, item, style)


def add_sidebar_item_paragraph(shape: Any, item: str, style: TextStyle) -> Any:
    """Render a teal bullet with black sidebar body text."""
    paragraph = add_paragraph(
        shape,
        "",
        size=style.body_size,
        color=TEXT_DARK,
        font=style.body_font,
    )
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("marL", "171450")
    p_pr.set("indent", "-91440")

    bullet_run = paragraph.add_run()
    bullet_run.text = "\u2022 "
    bullet_run.font.size = Pt(style.body_size)
    bullet_run.font.name = style.body_font
    bullet_run.font.color.rgb = rgb(AESG_TEAL)

    body_run = paragraph.add_run()
    body_run.text = clean_text(item)
    body_run.font.size = Pt(style.body_size)
    body_run.font.name = style.body_font
    body_run.font.color.rgb = rgb(TEXT_DARK)
    return paragraph


def add_experience_paragraph(shape: Any, item: str, style: TextStyle) -> Any:
    """Render one experience row with wrapped body lines aligned after its date."""
    prefix, separator, body = item.partition(EXPERIENCE_SEPARATOR)
    prefix = clean_text(prefix)
    body = clean_text(body) if separator else prefix
    if not separator:
        prefix = ""

    paragraph = add_paragraph(
        shape,
        "",
        size=style.body_size,
        color=TEXT_DARK,
        font=style.body_font,
    )
    p_pr = paragraph._p.get_or_add_pPr()
    if prefix:
        # Match the DOCX renderer: one shared tab aligns the first black role
        # character, while the matching hanging indent aligns every wrapped
        # continuation line to that same position.
        body_start = experience_body_start(prefix, style)
        p_pr.set("marL", str(body_start))
        p_pr.set("indent", str(-body_start))
        p_pr.set("algn", "l")
        for existing_tabs in list(p_pr.findall(qn("a:tabLst"))):
            p_pr.remove(existing_tabs)
        tab_list = OxmlElement("a:tabLst")
        tab = OxmlElement("a:tab")
        tab.set("pos", str(body_start))
        tab.set("algn", "l")
        tab_list.append(tab)
        p_pr.append(tab_list)

        prefix_run = paragraph.add_run()
        prefix_run.text = f"{prefix}\t"
        prefix_run.font.size = Pt(style.body_size)
        prefix_run.font.name = style.body_font
        prefix_run.font.color.rgb = rgb(AESG_TEAL)

        body_run = paragraph.add_run()
        body_run.text = body
        body_run.font.size = Pt(style.body_size)
        body_run.font.name = style.body_font
        body_run.font.color.rgb = rgb(TEXT_DARK)
    else:
        body_run = paragraph.add_run()
        body_run.text = body
        body_run.font.size = Pt(style.body_size)
        body_run.font.name = style.body_font
        body_run.font.color.rgb = rgb(TEXT_DARK)
    return paragraph


def add_block_items_text(shape: Any, blocks: list[tuple[str, list[str], bool]], style: TextStyle) -> None:
    clear_shape(shape)
    first_rendered = True
    for heading, items, show_heading in blocks:
        if not items:
            continue
        if show_heading:
            if not first_rendered:
                add_paragraph(shape, "", size=style.body_size, color=AESG_TEAL, font=style.body_font, space_after=2)
            add_paragraph(
                shape,
                heading.upper(),
                size=max(7.0, style.heading_size - 0.6),
                bold=True,
                color=AESG_TEAL,
                font=style.title_font,
                space_after=2,
            )
        for item in items:
            add_sidebar_item_paragraph(shape, item, style)
        first_rendered = False


def section_item_height_total(items: list[str], shape: Any, style: TextStyle) -> float:
    return sum(item_height(item, shape, style) for item in items)


def block_height_total(blocks: list[tuple[str, list[str], bool]], shape: Any, style: TextStyle) -> float:
    height = 0.0
    first_rendered = True
    for _heading, items, show_heading in blocks:
        if not items:
            continue
        if show_heading:
            if not first_rendered:
                height += style.body_size * 0.8
            height += max(7.0, style.heading_size - 0.6) * 1.28 + 3
        for item in items:
            height += item_height(item, shape, style)
        first_rendered = False
    return height


def fitted_style_for_height(
    shape: Any,
    style: TextStyle,
    estimated_height: float,
    *,
    min_body: float = 4.7,
    min_heading: float = 5.8,
) -> TextStyle:
    capacity = available_height_pt(shape, padding_pt=1)
    if estimated_height <= capacity:
        return style
    scale = max(0.58, min(1.0, capacity / max(estimated_height, 1.0)))
    return replace(
        style,
        body_size=max(min_body, style.body_size * scale * 0.96),
        heading_size=max(min_heading, style.heading_size * scale * 0.96),
    )


def chars_per_line(width_emu: int, size_pt: float, factor: float = 108.0) -> int:
    width_emu = width_emu or EMU_PER_INCH
    width_inches = max(0.6, width_emu / EMU_PER_INCH)
    return max(12, int(width_inches * factor / max(size_pt, 1.0)))


def wrapped_lines(text: str, width_emu: int, size_pt: float, factor: float = 108.0) -> int:
    text = clean_text(text)
    if not text:
        return 1
    per_line = chars_per_line(width_emu, size_pt, factor)
    total = 0
    for raw in text.splitlines() or [""]:
        total += max(1, math.ceil(len(raw) / per_line))
    return total


def paragraph_height(text: str, width_emu: int, size_pt: float, *, factor: float = 108.0, leading: float = 1.18) -> float:
    return wrapped_lines(text, width_emu, size_pt, factor) * size_pt * leading


def verdana_paragraph_height(text: str, width_emu: int, size_pt: float, *, bold: bool = False, leading: float = 1.14) -> float:
    """Measure portrait text with Verdana's real word widths.

    Character-count estimates are adequate for the compact landscape cards,
    but they overestimate the line capacity of the portrait template.  That
    caused a project to be accepted into a box that PowerPoint subsequently
    rendered past the bottom edge.
    """
    text = clean_text(text)
    if not text:
        return 0.0
    try:
        from PIL import ImageFont

        font_path = (
            "/System/Library/Fonts/Supplemental/Verdana Bold.ttf"
            if bold
            else "/System/Library/Fonts/Supplemental/Verdana.ttf"
        )
        font = ImageFont.truetype(font_path, max(1, round(size_pt * 96 / 72)))
        width_px = max(1.0, width_emu / EMU_PER_INCH * 96)
        lines = 0
        for raw_line in text.splitlines() or [text]:
            words = raw_line.split()
            if not words:
                lines += 1
                continue
            current = ""
            for word in words:
                candidate = f"{current} {word}".strip()
                if current and font.getlength(candidate) > width_px:
                    lines += 1
                    current = word
                else:
                    current = candidate
            if current:
                lines += 1
        return max(1, lines) * size_pt * leading
    except (OSError, AttributeError):
        return paragraph_height(text, width_emu, size_pt, factor=108.0, leading=leading)


def project_paragraph_height(
    text: str,
    width_emu: int,
    size_pt: float,
    style: TextStyle,
    *,
    bold: bool = False,
) -> float:
    """Return one conservative, shared project-text measurement.

    Project splitting and final box sizing must use the same measurement.
    Previously landscape split with a tight character estimate but then packed
    its final box with a much more optimistic estimate.  That disagreement
    allowed a project to be accepted into a box that could clip in PowerPoint.
    Portrait has a local Verdana measurement; landscape intentionally retains
    its source-template Ubuntu / Roboto Slab typography and uses a conservative
    line estimate because Ubuntu is supplied by the Office environment rather
    than this render host.
    """
    if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
        return verdana_paragraph_height(
            text, width_emu, size_pt, bold=bold, leading=1.14
        )
    # Landscape cards are compact and their body is Ubuntu.  A factor close to
    # the generic Office width avoids relying on placeholder-era dimensions;
    # it errs toward a continuation rather than clipping content.
    return paragraph_height(
        text,
        width_emu,
        size_pt,
        factor=104.0 if bold else 106.0,
        leading=1.16,
    )


def available_height_pt(shape: Any, padding_pt: float = 3.0) -> float:
    height = getattr(shape, "height", None) or EMU_PER_INCH
    return max(8.0, (height / 12700) - padding_pt)


def section_entry_height(
    heading: str,
    items: list[str],
    shape: Any,
    style: TextStyle,
    include_heading: bool,
    *,
    factor: float = 98.0,
) -> float:
    height = 0.0
    if include_heading:
        height += style.heading_size * 1.28 + 2
    for item in items:
        height += item_height(item, shape, style, factor=factor)
    return height


def split_long_item(
    item: str,
    shape: Any,
    style: TextStyle,
    include_heading: bool,
    *,
    factor: float = 98.0,
) -> list[str]:
    limit = max(40, int(chars_per_line(shape.width, style.body_size, factor=factor) * 5))
    words = item.split()
    chunks: list[str] = []
    current: list[str] = []
    for word in words:
        test = " ".join(current + [word])
        if current and len(test) > limit:
            chunks.append(" ".join(current))
            current = [word]
        else:
            current.append(word)
    if current:
        chunks.append(" ".join(current))
    return chunks or [item]


def item_height(item: str, shape: Any, style: TextStyle, *, factor: float = 98.0) -> float:
    """Measure a sidebar item using the same geometry as its renderer.

    Experience bodies use a hanging indent after the teal date.  Measuring
    them against the full text-box width underestimates wrapped entries and
    lets the following project heading overlap their final line.
    """
    prefix, separator, body = item.partition(EXPERIENCE_SEPARATOR)
    if separator and clean_text(prefix):
        body_width = max(
            pt_to_emu(style.body_size * 8),
            shape.width - experience_body_start(clean_text(prefix), style),
        )
        return paragraph_height(
            clean_text(body),
            body_width,
            style.body_size,
            factor=factor,
            leading=1.16,
        ) + 1.0
    return paragraph_height(item, shape.width, style.body_size, factor=factor, leading=1.16) + 1.0


def pack_items_into_shape(items: list[str], shape: Any, style: TextStyle) -> tuple[list[str], list[str]]:
    capacity = available_height_pt(shape)
    used = 0.0
    packed: list[str] = []
    for item_index, item in enumerate(items):
        pieces = split_long_item(item, shape, style, include_heading=False)
        for piece_index, piece in enumerate(pieces):
            height = item_height(piece, shape, style)
            if used + height <= capacity or not packed:
                packed.append(piece)
                used += height
            else:
                return packed, pieces[piece_index:] + items[item_index + 1 :]
    return packed, []


def pack_blocks_into_shape(
    blocks: list[tuple[str, list[str], bool]],
    shape: Any,
    style: TextStyle,
) -> tuple[list[tuple[str, list[str], bool]], list[tuple[str, list[str], bool]]]:
    capacity = available_height_pt(shape)
    used = 0.0
    packed_blocks: list[tuple[str, list[str], bool]] = []
    for block_index, (heading, items, show_heading) in enumerate(blocks):
        full_height = block_height_total([(heading, items, show_heading)], shape, style)
        # Qualifications, Education, and Memberships are semantic blocks. If
        # a complete block fits on an empty slide, defer it intact rather than
        # leaving a heading or the first credential behind in the remainder.
        if full_height <= capacity and used + full_height > capacity and packed_blocks:
            return packed_blocks, blocks[block_index:]
        packed_items: list[str] = []
        heading_height = max(7.0, style.heading_size - 0.6) * 1.28 + 3 if show_heading else 0.0
        for item_index, item in enumerate(items):
            pieces = split_long_item(item, shape, style, include_heading=show_heading and not packed_items)
            for piece_index, piece in enumerate(pieces):
                height = item_height(piece, shape, style)
                if show_heading and not packed_items:
                    height += heading_height
                if used + height <= capacity or not packed_blocks and not packed_items:
                    packed_items.append(piece)
                    used += height
                else:
                    if packed_items:
                        packed_blocks.append((heading, packed_items, show_heading))
                    remaining_block = (heading, pieces[piece_index:] + items[item_index + 1 :], show_heading)
                    return packed_blocks, [remaining_block] + blocks[block_index + 1 :]
        if packed_items:
            packed_blocks.append((heading, packed_items, show_heading))
    return packed_blocks, []


def pack_sections_into_shape(
    sections: list[tuple[str, list[str], bool]],
    shape: Any,
    style: TextStyle,
    *,
    density_factor: float = 98.0,
    keep_sections_together: bool = True,
) -> tuple[list[tuple[str, list[str], bool]], list[tuple[str, list[str], bool]]]:
    capacity = available_height_pt(shape)
    used = 0.0
    packed: list[tuple[str, list[str], bool]] = []
    remaining: list[tuple[str, list[str], bool]] = []

    for section_index, (heading, items, continued) in enumerate(sections):
        include_heading = not continued
        # Landscape keeps sidebar sections as semantic units. Portrait should
        # use every available line before continuing on the next slide, just
        # like the flowing portrait DOCX row.
        atomic_section = keep_sections_together
        full_height = section_entry_height(
            heading,
            items,
            shape,
            style,
            include_heading,
            factor=density_factor,
        )
        if atomic_section and full_height <= capacity:
            if used + full_height <= capacity:
                packed.append((heading, list(items), continued))
                used += full_height
                continue
            if packed:
                remaining.append((heading, list(items), continued))
                remaining.extend(sections[section_index + 1 :])
                return packed, remaining

        bucket: list[str] = []
        for item_index, item in enumerate(items):
            include_item_heading = include_heading and not bucket
            full_item_height = section_entry_height(
                heading,
                [item],
                shape,
                style,
                include_item_heading,
                factor=density_factor,
            )
            # Keep a complete sidebar item intact whenever it can fit on an
            # otherwise empty slide. Split only an intrinsically oversize
            # item, never merely because the current slide is nearly full.
            chunks = (
                [item]
                if full_item_height <= capacity
                else split_long_item(
                    item,
                    shape,
                    style,
                    include_item_heading,
                    factor=density_factor,
                )
            )
            for part_index, piece in enumerate(chunks):
                height = section_entry_height(
                    heading,
                    [piece],
                    shape,
                    style,
                    include_heading and not bucket,
                    factor=density_factor,
                )
                if used + height <= capacity or not packed and not bucket:
                    bucket.append(piece)
                    used += height
                else:
                    if bucket:
                        packed.append((heading, bucket, continued))
                    leftover = chunks[part_index:] + items[item_index + 1 :]
                    # Do not suppress the heading on the next slide when no
                    # item from this section was placed on the current slide.
                    remaining.append((heading, leftover, continued or bool(bucket)))
                    remaining.extend(sections[section_index + 1 :])
                    return packed, remaining
        if bucket:
            packed.append((heading, bucket, continued))
    return packed, remaining


def pack_sections(
    sections: list[tuple[str, list[str], bool]],
    shapes: list[Any],
    style: TextStyle,
    *,
    keep_sections_together: bool = True,
) -> list[list[tuple[str, list[str], bool]]]:
    packed_by_shape: list[list[tuple[str, list[str], bool]]] = []
    remaining = sections
    for shape in shapes:
        packed, remaining = pack_sections_into_shape(
            remaining,
            shape,
            style,
            keep_sections_together=keep_sections_together,
        )
        packed_by_shape.append(packed)
        if not remaining:
            break
    while len(packed_by_shape) < len(shapes):
        packed_by_shape.append([])
    return packed_by_shape


def estimate_project_fixed_height(project: dict[str, Any], shape: Any, style: TextStyle) -> float:
    title = clean_text(project.get("name") or "Project")
    client = project_client_display(project)
    duration = clean_text(project.get("duration_line") or project.get("duration") or "")
    title_line = title
    height = paragraph_height(title_line, shape.width, style.project_title_size, factor=102.0, leading=1.18) + 2
    if client:
        height += paragraph_height(client, shape.width, style.project_body_size, factor=104.0, leading=1.16) + 1
    if duration:
        height += paragraph_height(duration, shape.width, style.project_body_size, factor=104.0, leading=1.16) + 2
    return height + 2


def body_capacity_chars(project: dict[str, Any], shape: Any, style: TextStyle) -> int:
    capacity = available_height_pt(shape, padding_pt=4)
    fixed = estimate_project_fixed_height(project, shape, style)
    remaining_pt = max(style.project_body_size * 2, capacity - fixed)
    lines = max(1, int(remaining_pt / (style.project_body_size * 1.16)))
    per_line = chars_per_line(shape.width, style.project_body_size, factor=104.0)
    return max(45, int(lines * per_line * 1.05))


def split_body(text: str, limit: int) -> list[str]:
    text = clean_text(text)
    if not text:
        return [""]
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: list[str] = []
    current = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        candidate = f"{current} {sentence}".strip()
        if current and len(candidate) > limit:
            chunks.append(current)
            current = sentence
        elif len(sentence) > limit:
            if current:
                chunks.append(current)
                current = ""
            words = sentence.split()
            word_chunk: list[str] = []
            for word in words:
                candidate_word = " ".join(word_chunk + [word])
                if word_chunk and len(candidate_word) > limit:
                    chunks.append(" ".join(word_chunk))
                    word_chunk = [word]
                else:
                    word_chunk.append(word)
            current = " ".join(word_chunk)
        else:
            current = candidate
    if current:
        chunks.append(current)
    return chunks or [""]


def split_project_body_to_height(project: dict[str, Any], shape: Any, style: TextStyle) -> list[str]:
    """Split a project only when its final rendered block cannot fit.

    Portrait cards use Verdana, whose proportional word widths differ too much
    from the legacy character-count approximation to use it as a pagination
    decision.  Build complete sentence groups and check each prospective group
    against the *actual* rendered project-height estimator instead.  This
    keeps normal projects intact while making an intrinsically tall project a
    clean continuation rather than allowing PowerPoint to clip it.
    """
    body = clean_text(project.get("description") or "")
    if not body:
        return [""]

    capacity = max(pt_to_emu(style.project_body_size * 4), shape.height - pt_to_emu(2))

    def fits(value: str) -> bool:
        candidate = ProjectChunk(project=project, body=value, part=1, parts=1)
        return estimate_project_chunk_height(candidate, shape, style) <= capacity

    sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", body) if part.strip()]
    result: list[str] = []
    current = ""
    for sentence in sentences:
        candidate = f"{current} {sentence}".strip()
        if current and not fits(candidate):
            result.append(current)
            current = ""
        if not current and not fits(sentence):
            # A single long sentence must be broken at a word boundary.  The
            # same height test determines each fragment, not a guessed number
            # of characters.
            word_group = ""
            for word in sentence.split():
                candidate_word = f"{word_group} {word}".strip()
                if word_group and not fits(candidate_word):
                    result.append(word_group)
                    word_group = word
                else:
                    word_group = candidate_word
            current = word_group
        else:
            current = sentence if not current else candidate
    if current:
        result.append(current)
    return result or [""]


def split_project_chunk_to_height(
    chunk: ProjectChunk,
    shape: Any,
    capacity: int,
    style: TextStyle,
) -> list[ProjectChunk]:
    """Break one already-created project chunk at sentence/word boundaries.

    This is the last safety gate after a front-slide Experience block has
    pushed the project lane down.  Earlier chunking has no way to know that
    final available height, so this function deliberately lives in the
    renderer and is used only when a whole project block otherwise cannot fit.
    """
    body = clean_text(chunk.body)
    if not body:
        return [chunk]

    def fits(value: str) -> bool:
        candidate = ProjectChunk(
            project=chunk.project,
            body=value,
            part=chunk.part,
            parts=chunk.parts,
        )
        return estimate_project_chunk_height(candidate, shape, style) <= capacity

    sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", body) if part.strip()]
    bodies: list[str] = []
    current = ""
    for sentence in sentences:
        candidate = f"{current} {sentence}".strip()
        if current and not fits(candidate):
            bodies.append(current)
            current = ""
        if not current and not fits(sentence):
            words = ""
            for word in sentence.split():
                candidate_word = f"{words} {word}".strip()
                if words and not fits(candidate_word):
                    bodies.append(words)
                    words = word
                else:
                    words = candidate_word
            current = words
        else:
            current = sentence if not current else candidate
    if current:
        bodies.append(current)
    if len(bodies) <= 1:
        return [chunk]

    first_part = chunk.part
    total_parts = first_part + len(bodies) - 1
    return [
        ProjectChunk(chunk.project, body, first_part + index, total_parts)
        for index, body in enumerate(bodies)
    ]


def project_chunks_for_shapes(
    projects: list[dict[str, Any]],
    prototype_shape: Any,
    style: TextStyle,
    *,
    use_teaser: bool = False,
) -> list[ProjectChunk]:
    chunks: list[ProjectChunk] = []
    for project in projects:
        if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
            # The front-page placeholders are intentionally short. Splitting
            # against those boxes creates artificial "continued 2/6" cards
            # across columns. Start with one atomic project and let a
            # project-only continuation page split it only if the complete
            # project is taller than that entire page.
            description = (
                project.get("description_teaser")
                if use_teaser
                else project.get("description")
            )
            body_chunks = [
                clean_text(
                    description
                    or project.get("description")
                    or project.get("description_teaser")
                    or ""
                )
            ]
        else:
            limit = body_capacity_chars(project, prototype_shape, style)
            body_chunks = split_body(project.get("description") or "", limit)
        total = len(body_chunks)
        for index, body in enumerate(body_chunks, start=1):
            chunks.append(ProjectChunk(project=project, body=body, part=index, parts=total))
    return chunks


def project_identity(chunk: ProjectChunk) -> tuple[Any, str, str, str]:
    project = chunk.project
    return (
        project.get("source_index", project.get("rank")),
        clean_text(project.get("name") or ""),
        project_client_display(project),
        clean_text(project.get("duration_line") or project.get("duration") or ""),
    )


def merge_project_chunks_for_shape(
    chunks: list[ProjectChunk],
    prototype_shape: Any,
    style: TextStyle,
    max_height: int,
) -> list[ProjectChunk]:
    merged: list[ProjectChunk] = []
    index = 0
    while index < len(chunks):
        group = [chunks[index]]
        identity = project_identity(chunks[index])
        index += 1
        while index < len(chunks) and project_identity(chunks[index]) == identity:
            group.append(chunks[index])
            index += 1

        if len(group) == 1:
            only = group[0]
            if only.part > 1:
                merged.append(ProjectChunk(project=only.project, body=only.body, part=1, parts=1))
            else:
                merged.append(only)
            continue

        combined = ProjectChunk(
            project=group[0].project,
            body=clean_text(" ".join(chunk.body for chunk in group if chunk.body)),
            part=1,
            parts=1,
        )
        if estimate_project_chunk_height(combined, prototype_shape, style) <= max_height:
            merged.append(combined)
        else:
            merged.extend(group)
    return merged


def find_textbox_shapes(slide: Any, indexes: Iterable[int]) -> list[Any]:
    return [slide.shapes[index] for index in indexes if index < len(slide.shapes)]


def duplicate_from_template(prs: Presentation, template_index: int) -> Any:
    source = prs.slides[template_index]
    return duplicate_from_snapshot(prs, source.slide_layout, [copy.deepcopy(shape.element) for shape in source.shapes])


def duplicate_from_snapshot(prs: Presentation, layout: Any, elements: list[Any]) -> Any:
    """Duplicate an immutable template snapshot, never a populated slide."""
    destination = prs.slides.add_slide(layout)
    for shape in list(destination.shapes):
        shape.element.getparent().remove(shape.element)
    for element in elements:
        destination.shapes._spTree.insert_element_before(copy.deepcopy(element), "p:extLst")
    return destination


def remove_unwanted_slides(prs: Presentation, keep_indexes: set[int]) -> None:
    for index in range(len(prs.slides) - 1, -1, -1):
        if index not in keep_indexes:
            delete_pptx_slide(prs, index)


def photo_path(
    cv: dict[str, Any],
    tmpdir: Path,
    box_w: int,
    box_h: int,
    background_color: str = AESG_TEAL,
) -> Path | None:
    photo = cv.get("profile_photo") or {}
    encoded = photo.get("base64")
    if not encoded:
        return None
    try:
        data = base64.b64decode(encoded)
    except Exception:
        return None
    source = Image.open(io.BytesIO(data)).convert("RGB")
    target_w = max(100, int((box_w / EMU_PER_INCH) * 240))
    target_h = max(100, int((box_h / EMU_PER_INCH) * 240))
    fitted = ImageOps.contain(source, (target_w, target_h), method=Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (target_w, target_h), f"#{background_color}")
    offset = ((target_w - fitted.width) // 2, (target_h - fitted.height) // 2)
    canvas.paste(fitted, offset)
    output = tmpdir / f"{slugify(cv.get('name') or 'photo')}.png"
    canvas.save(output, "PNG")
    return output


def replace_photo(
    shape: Any,
    cv: dict[str, Any],
    tmpdir: Path,
    background_color: str = AESG_TEAL,
) -> None:
    path = photo_path(cv, tmpdir, shape.width, shape.height, background_color)
    if not path:
        clear_shape(shape)
        return
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    clear_shape(shape)
    shape.part.slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def profile_parts(cv: dict[str, Any]) -> tuple[str, str, str]:
    name = clean_text(cv.get("name") or cv.get("metadata", {}).get("name") or "Unnamed CV")
    first = clean_text(cv.get("first_name") or cv.get("metadata", {}).get("first_name") or "")
    last = clean_text(cv.get("last_name") or cv.get("metadata", {}).get("last_name") or "")
    if not first:
        parts = name.split(maxsplit=1)
        first = parts[0] if parts else name
        last = parts[1] if len(parts) > 1 else ""
    role = clean_text(cv.get("role") or cv.get("metadata", {}).get("role") or "")
    return first, last, role


def section_data(cv: dict[str, Any]) -> list[tuple[str, list[str], bool]]:
    sections = [
        ("Experience", experience_line_items(cv.get("experience"))),
        ("Key Expertise", line_items(cv.get("key_expertise"))),
        ("Qualifications", line_items(cv.get("qualifications"))),
        ("Education", line_items(cv.get("education"))),
        ("Memberships", line_items(cv.get("memberships"))),
    ]
    return [(heading, items, False) for heading, items in sections if items]


def selected_projects(cv: dict[str, Any]) -> list[dict[str, Any]]:
    projects = cv.get("selected_projects") or []
    return [project for project in projects if clean_text(project.get("name") or "")]


def standard_remainders(cv: dict[str, Any]) -> dict[str, list[str]]:
    return {
        "experience": experience_line_items(cv.get("experience")),
        "key": line_items(cv.get("key_expertise")),
        "qualifications": line_items(cv.get("qualifications")),
        "education": line_items(cv.get("education")),
        "memberships": line_items(cv.get("memberships")),
    }


def qualification_blocks(remainders: dict[str, list[str]]) -> list[tuple[str, list[str], bool]]:
    blocks: list[tuple[str, list[str], bool]] = []
    if remainders.get("qualifications"):
        blocks.append(("Qualifications", remainders["qualifications"], True))
    if remainders.get("education"):
        blocks.append(("Education", remainders["education"], True))
    if remainders.get("memberships"):
        blocks.append(("Memberships", remainders["memberships"], True))
    return blocks


def apply_qualification_blocks(remainders: dict[str, list[str]], remaining_blocks: list[tuple[str, list[str], bool]]) -> None:
    remainders["qualifications"] = []
    remainders["education"] = []
    remainders["memberships"] = []
    for heading, items, _show_heading in remaining_blocks:
        key = heading.lower()
        if key == "qualifications":
            remainders["qualifications"] = items
        elif key == "education":
            remainders["education"] = items
        elif key == "memberships":
            remainders["memberships"] = items


def section_continuations_from_remainders(remainders: dict[str, list[str]]) -> list[tuple[str, list[str], bool]]:
    sections: list[tuple[str, list[str], bool]] = []
    if remainders.get("experience"):
        sections.append(("Experience", remainders["experience"], True))
    if remainders.get("key"):
        sections.append(("Key Expertise", remainders["key"], True))
    if remainders.get("qualifications"):
        sections.append(("Qualifications", remainders["qualifications"], True))
    if remainders.get("education"):
        sections.append(("Education", remainders["education"], True))
    if remainders.get("memberships"):
        sections.append(("Memberships", remainders["memberships"], True))
    return sections


def remainders_have_content(remainders: dict[str, list[str]]) -> bool:
    return any(items for items in remainders.values())


def blank_shape(shape: Any) -> None:
    clear_shape(shape)


def cover_static_project_heading(slide: Any, shape: Any) -> None:
    # Static labels are suppressed directly on their layouts. Keeping this as
    # a no-op avoids adding white rectangles that Google Slides styles as cards.
    return


def set_shape_geometry(shape: Any, left: int, top: int, width: int, height: int) -> None:
    xfrm = shape.element.get_or_add_xfrm()
    off = xfrm.off
    if off is None:
        off = OxmlElement("a:off")
        xfrm.insert(0, off)
    ext = xfrm.ext
    if ext is None:
        ext = OxmlElement("a:ext")
        xfrm.insert(1, ext)
    off.set("x", str(int(left)))
    off.set("y", str(int(top)))
    ext.set("cx", str(int(width)))
    ext.set("cy", str(int(height)))


def send_shape_to_back(shape: Any) -> None:
    parent = shape.element.getparent()
    if parent is None:
        return
    parent.remove(shape.element)
    parent.insert(2, shape.element)


def cover_layout_labels(slide: Any, labels: set[str] = LAYOUT_LABELS) -> None:
    return


def suppress_layout_labels(prs: Presentation, labels: set[str] = LAYOUT_LABELS) -> None:
    """Remove variable labels from layouts so generated headings can move."""
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if clean_text(shape.text).upper() in labels:
                shape.text_frame.clear()


def suppress_portrait_vertical_layout_dividers(prs: Presentation) -> None:
    """Replace inherited portrait dividers with consistently sized slide rules."""
    for layout in prs.slide_layouts:
        for shape in list(layout.shapes):
            is_vertical_line = (
                getattr(shape, "shape_type", None) == 9
                and shape.height > int(0.75 * EMU_PER_INCH)
                and abs(shape.width) < 30000
            )
            if not is_vertical_line:
                continue
            parent = shape.element.getparent()
            if parent is not None:
                parent.remove(shape.element)


def remove_empty_placeholders(prs: Presentation) -> None:
    """Delete unused placeholders that Google Slides exposes as empty boxes."""
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not getattr(shape, "is_placeholder", False):
                continue
            if getattr(shape, "has_text_frame", False) and clean_text(shape.text):
                continue
            parent = shape.element.getparent()
            if parent is not None:
                parent.remove(shape.element)


def layout_svg_blob(shape: Any) -> bytes | None:
    """Return the SVG-only image used by the original template layout."""
    relationship_id = None
    for element in shape.element.iter():
        if element.tag.endswith("svgBlip"):
            for key, value in element.attrib.items():
                if key.endswith("}embed"):
                    relationship_id = value
                    break
        if relationship_id:
            break
    if not relationship_id:
        return None
    try:
        return shape.part.related_part(relationship_id).blob
    except Exception:
        return None


def promote_svg_layout_logo(
    slide: Any,
    tmpdir: Path,
    slide_number: int,
) -> bool:
    """Make the template's SVG logo portable to Google Slides.

    The source template uses an SVG-only layout picture with no PNG fallback.
    We rasterize that exact SVG and place it at the exact same coordinates.
    """
    promoted = False
    for picture_index, shape in enumerate(slide.slide_layout.shapes):
        if getattr(shape, "shape_type", None) != 13:
            continue
        blob = layout_svg_blob(shape)
        if not blob:
            continue

        png_path = WHITE_LOGO_ASSET
        if not png_path.exists():
            # Keep a best-effort conversion path for custom templates, but do
            # not remove the original layout SVG unless a fallback was
            # successfully created.
            svg_path = tmpdir / f"layout_logo_{slide_number}_{picture_index}.svg"
            png_path = tmpdir / f"layout_logo_{slide_number}_{picture_index}.png"
            svg_path.write_bytes(blob)
            try:
                subprocess.run(
                    ["rsvg-convert", "-o", str(png_path), str(svg_path)],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
            except (OSError, subprocess.CalledProcessError):
                continue

        slide.shapes.add_picture(
            str(png_path),
            shape.left,
            shape.top,
            width=shape.width,
            height=shape.height,
        )
        promoted = True
    return promoted


def finalize_presentation(prs: Presentation, tmpdir: Path) -> None:
    remove_empty_placeholders(prs)
    promoted_layout_ids: set[int] = set()
    for slide_number, slide in enumerate(prs.slides):
        if promote_svg_layout_logo(slide, tmpdir, slide_number):
            promoted_layout_ids.add(id(slide.slide_layout))
    # The PNG fallback is now on every affected slide, so remove the SVG-only
    # layout object that Google Slides imports as an outlined blank rectangle.
    for layout in prs.slide_layouts:
        if id(layout) not in promoted_layout_ids:
            continue
        for shape in list(layout.shapes):
            if getattr(shape, "shape_type", None) != 13 or not layout_svg_blob(shape):
                continue
            parent = shape.element.getparent()
            if parent is not None:
                parent.remove(shape.element)


def add_dynamic_heading(
    slide: Any,
    text: str,
    left: int,
    top: int,
    width: int,
    style: TextStyle,
    *,
    with_rule: bool = False,
) -> None:
    # Keep the heading text box at the full project area width so the text
    # never overflows; draw the underline only under the heading text itself.
    # The template uses a 1-inch-wide, 1.5 pt teal straight connector with no
    # shadow (effectRef idx="0"). Replicate that exact object style.
    TEMPLATE_LINE_WIDTH = 914400  # 1 inch
    TEMPLATE_LINE_WEIGHT = 19050  # 1.5 pt
    shape_height = pt_to_emu(style.heading_size * (1.45 if with_rule else 2.3)) + (19050 if with_rule else 0)
    shape = slide.shapes.add_textbox(left, max(0, top), width, shape_height)
    clear_shape(shape)
    paragraph = add_paragraph(
        shape,
        text.upper(),
        size=style.heading_size,
        bold=True,
        color=AESG_TEAL,
        font=style.title_font,
    )
    if with_rule:
        # Place the exact template-style line under the heading text.
        line_y = shape.top + shape.height
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            shape.left,
            line_y,
            shape.left + TEMPLATE_LINE_WIDTH,
            line_y,
        )
        line.line.color.rgb = RGBColor.from_string(AESG_TEAL)
        line.line.width = Pt(1.5)
        # Remove the inherited theme style so no shadow/3D effect is applied.
        style_el = line.element.find(qn("p:style"))
        if style_el is not None:
            line.element.remove(style_el)


def slide_has_project_heading(slide: Any) -> bool:
    """Return True if the slide already contains a dynamic 'Selected Projects' heading."""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = clean_text(shape.text)
            if text and "SELECTED PROJECTS" in text.upper():
                return True
    return False


def slide_layout_has_teal_banner(slide: Any) -> bool:
    """Return True if the slide's layout has a wide top banner strip."""
    for shape in slide.slide_layout.shapes:
        if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.top > int(0.1 * EMU_PER_INCH):
            continue
        if shape.left > int(0.1 * EMU_PER_INCH):
            continue
        if shape.width < int(7.0 * EMU_PER_INCH):
            continue
        if shape.height > int(2.0 * EMU_PER_INCH):
            continue
        if shape.height < int(0.5 * EMU_PER_INCH):
            continue
        return True
    return False


def add_project_area_heading(
    slide: Any,
    project_shapes: list[Any],
    style: TextStyle,
    *,
    heading_top: int | None = None,
    project_only_full_width: bool = False,
    with_rule: bool = False,
) -> None:
    filled = [shape for shape in project_shapes if shape_has_text(shape)]
    if not filled:
        return
    top = min(shape.top for shape in filled)
    left = min(shape.left for shape in filled)
    width = max(shape.left + shape.width for shape in filled) - left
    top = heading_top if heading_top is not None else top - pt_to_emu(style.heading_size * (2.05 if with_rule else 2.6)) - (19050 if with_rule else 0)
    # The portrait continuation template leaves a slightly larger optical
    # buffer below the compact AESG banner before the section heading. The
    # project boxes themselves already use the correct source-template
    # coordinates, so move only this generated heading.
    if project_only_full_width and style is PORTRAIT_STYLE and heading_top is None:
        top += int(0.12 * EMU_PER_INCH)
    if style is LANDSCAPE_STYLE:
        # Landscape continuation artwork uses a shallow top banner. Keep the
        # dynamic heading fully below it with a small optical buffer instead
        # of allowing the text box to begin on the banner edge.
        top = max(top, int(0.95 * EMU_PER_INCH))
    add_dynamic_heading(slide, "Selected Projects", left, top, width, style, with_rule=with_rule)


def fill_project_shapes(shapes: list[Any], chunks: list[ProjectChunk], style: TextStyle) -> tuple[int, list[Any]]:
    used = 0
    for index, shape in enumerate(shapes):
        if index < len(chunks):
            add_project_text(shape, chunks[index], style)
            used += 1
        else:
            blank_shape(shape)
    return used, shapes


def pt_to_emu(value: float) -> int:
    return int(value * 12700)


def estimate_project_chunk_height(chunk: ProjectChunk, shape: Any, style: TextStyle) -> int:
    project = chunk.project
    title = clean_text(project.get("name") or "Project")
    if chunk.parts > 1 and chunk.part > 1:
        title = f"{title} (continued {chunk.part}/{chunk.parts})"
    client = project_client_display(project)
    duration = clean_text(project.get("duration_line") or project.get("duration") or "")
    measure = lambda value, size, bold=False: project_paragraph_height(
        value, shape.width, size, style, bold=bold
    )
    height_pt = measure(title, style.project_title_size, True) + 4
    if client and chunk.part == 1:
        height_pt += measure(client, style.project_body_size) + 1
    if duration and chunk.part == 1:
        height_pt += measure(duration, style.project_body_size) + 3
    if chunk.body:
        height_pt += measure(chunk.body, style.project_body_size) + 3
    return pt_to_emu(height_pt)


def paragraph_text(paragraph: Any) -> str:
    return clean_text("".join(run.text for run in paragraph.runs) or paragraph.text)


def paragraph_font_size(paragraph: Any, default_size: float) -> float:
    sizes: list[float] = []
    if paragraph.font.size:
        sizes.append(paragraph.font.size.pt)
    for run in paragraph.runs:
        if run.font.size:
            sizes.append(run.font.size.pt)
    return max(sizes) if sizes else default_size


def paragraph_has_color(paragraph: Any, value: str) -> bool:
    for run in paragraph.runs:
        try:
            color = run.font.color.rgb
        except (AttributeError, TypeError):
            continue
        if color is not None and str(color) == value:
            return True
    return False


def is_rendered_experience_paragraph(paragraph: Any) -> bool:
    """Recognise the teal-date / dark-body experience paragraph structure."""
    return paragraph_has_color(paragraph, AESG_TEAL) and paragraph_has_color(
        paragraph, TEXT_DARK
    )


def rendered_experience_height(paragraph: Any, shape: Any, style: TextStyle) -> float:
    """Return the actual line height of a rendered hanging-indent row."""
    prefix = ""
    body_parts: list[str] = []
    for run in paragraph.runs:
        run_text = clean_text(run.text)
        if not run_text:
            continue
        try:
            run_color = run.font.color.rgb
        except (AttributeError, TypeError):
            run_color = None
        if str(run_color) == AESG_TEAL and not body_parts:
            prefix += run_text
        else:
            body_parts.append(run_text)
    body_width = max(
        pt_to_emu(style.body_size * 8),
        shape.width - experience_body_start(prefix, style),
    )
    return paragraph_height(
        " ".join(body_parts),
        body_width,
        style.body_size,
        factor=98.0,
        leading=1.16,
    )


def shape_has_text(shape: Any) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and clean_text(shape.text))


def estimate_filled_shape_height(shape: Any, style: TextStyle) -> int:
    if not shape_has_text(shape):
        return 0
    height_pt = 0.0
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph_text(paragraph)
        if not text:
            continue
        size = paragraph_font_size(paragraph, style.body_size)
        if is_rendered_experience_paragraph(paragraph):
            height_pt += rendered_experience_height(paragraph, shape, style)
        else:
            factor = 102.0 if size >= style.heading_size else 98.0
            height_pt += paragraph_height(text, shape.width, size, factor=factor, leading=1.17)
        if paragraph.space_after:
            height_pt += paragraph.space_after.pt
    return pt_to_emu(height_pt + 8)


def estimate_project_filled_shape_height(shape: Any, style: TextStyle) -> int:
    """Measure project text using Office-like character density.

    Generic CV bullets wrap more aggressively than justified project prose.
    A dedicated estimate avoids preserving placeholder-sized gaps between
    project blocks while retaining enough height for titles and role lines.
    """
    if not shape_has_text(shape):
        return 0
    height_pt = 0.0
    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph_text(paragraph)]
    for index, paragraph in enumerate(paragraphs):
        text = paragraph_text(paragraph)
        size = paragraph_font_size(paragraph, style.project_body_size)
        is_title = index == 0 or paragraph_has_color(paragraph, AESG_TEAL)
        height_pt += project_paragraph_height(
            text,
            shape.width,
            size,
            style,
            bold=is_title,
        )
        if paragraph.space_after:
            height_pt += paragraph.space_after.pt
    # Text frames have zero internal margins. In portrait, reserve one
    # additional rendered body line below the final paragraph. PowerPoint and
    # Google Slides can make a final justified line fractionally taller than
    # the Pillow measurement; without this safety band, a project that exactly
    # fills its calculated box is visibly clipped at the bottom. The allowance
    # is included in packing, so it creates a continuation slide when needed
    # rather than silently losing the final lines.
    bottom_safety = style.project_body_size * (
        1.35 if style is PORTRAIT_STYLE else 1.15
    )
    return pt_to_emu(height_pt + bottom_safety)


def estimate_compact_project_chunk_height(chunk: ProjectChunk, shape: Any, style: TextStyle) -> int:
    """Measure a landscape project block without placeholder-era padding."""
    project = chunk.project
    title = clean_text(project.get("name") or "Project")
    client = project_client_display(project)
    duration = clean_text(project.get("duration_line") or project.get("duration") or "")
    title_line = title
    height_pt = paragraph_height(
        title_line,
        shape.width,
        style.project_title_size,
        factor=125.0,
        leading=1.13,
    ) + 2
    if client and chunk.part == 1:
        height_pt += paragraph_height(
            client,
            shape.width,
            style.project_body_size,
            factor=145.0,
            leading=1.12,
        ) + 1
    if duration and chunk.part == 1:
        height_pt += paragraph_height(
            duration,
            shape.width,
            style.project_body_size,
            factor=145.0,
            leading=1.12,
        ) + 2
    if chunk.body:
        height_pt += paragraph_height(
            chunk.body,
            shape.width,
            style.project_body_size,
            factor=150.0,
            leading=1.12,
        )
    return pt_to_emu(height_pt + 2)


def estimate_compact_project_filled_height(shape: Any, style: TextStyle) -> int:
    """Re-measure a rendered landscape project using its actual paragraphs."""
    if not shape_has_text(shape):
        return 0
    height_pt = 0.0
    paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph_text(paragraph)]
    for index, paragraph in enumerate(paragraphs):
        text = paragraph_text(paragraph)
        size = paragraph_font_size(paragraph, style.project_body_size)
        if index == 0:
            factor, leading = 125.0, 1.13
        elif text.lower().startswith("role:"):
            factor, leading = 145.0, 1.12
        else:
            factor, leading = 150.0, 1.12
        height_pt += paragraph_height(text, shape.width, size, factor=factor, leading=leading)
        if paragraph.space_after:
            height_pt += paragraph.space_after.pt
    return pt_to_emu(height_pt + 2)


def measured_experience_body_start(prefix: str, style: TextStyle) -> int:
    """Measure the hanging-tab position required by one date prefix.

    PowerPoint stores paragraph offsets in EMUs and does not automatically
    align wrapped runs after a coloured prefix.  This mirrors the landscape
    DOCX measurement, but deliberately returns only the width of this one
    prefix.  ``configure_experience_body_start`` promotes the widest value to
    the shared CV-wide tab stop used by every experience row.
    """
    # Measure the visible date (including its following space) in the same
    # Verdana face as the output.  A character-count approximation was close
    # enough for short rows, but made wrapped body lines land a space to the
    # left or right of their first body line.
    if style.body_font.casefold() == "verdana":
        try:
            from PIL import ImageFont

            font_size_px = max(1, round(style.body_size * 96 / 72))
            font = ImageFont.truetype(
                "/System/Library/Fonts/Supplemental/Verdana.ttf",
                font_size_px,
            )
            width_pt = font.getlength(f"{prefix.strip()} ") * 72 / 96
            # A small fixed gutter accounts for PowerPoint's run metrics and
            # makes a wrapped body line land directly under the first black
            # body character rather than one character to its left.
            return pt_to_emu(width_pt + 1.5)
        except (OSError, AttributeError):
            pass

    # Fallback for environments without Verdana.
    width_units = 0.0
    for char in prefix.strip():
        if char.isspace():
            width_units += 0.28
        elif char in ".,:;()":
            width_units += 0.30
        elif char in "-\u2013\u2014/":
            width_units += 0.52
        elif char.isdigit():
            width_units += 0.55
        else:
            width_units += 0.53
    prefix_width = width_units * style.body_size
    clearance = 2.0
    return pt_to_emu(prefix_width + clearance)


def configure_experience_body_start(style: TextStyle, cv: dict[str, Any]) -> int:
    """Set one experience body start from the widest date in this CV.

    The landscape DOCX renderer measures all date labels before laying out
    any row.  PPTX must do the same so that the first black role character and
    every wrapped continuation line share one vertical plane.  Storing the
    value on the existing orientation style also makes all packing and
    overflow calculations consume the identical width.
    """
    prefixes = [
        clean_text(item.partition(EXPERIENCE_SEPARATOR)[0])
        for item in experience_line_items(cv.get("experience"))
        if EXPERIENCE_SEPARATOR in item
        and clean_text(item.partition(EXPERIENCE_SEPARATOR)[0])
    ]
    style.experience_body_start = max(
        (
            measured_experience_body_start(prefix, style)
            for prefix in prefixes
        ),
        default=0,
    )
    return style.experience_body_start


def experience_body_start(prefix: str, style: TextStyle) -> int:
    """Return the CV-wide hanging-tab position for an experience row."""
    measured = measured_experience_body_start(prefix, style)
    return max(style.experience_body_start, measured)


def compact_project_gap(style: TextStyle) -> int:
    """Use approximately one body-text line between project blocks."""
    return pt_to_emu(max(4.0, style.project_body_size * 1.0))


def compact_project_min_height(style: TextStyle, original_height: int) -> int:
    """Avoid placeholder-sized cards while retaining a small editable box."""
    return min(original_height, pt_to_emu(max(7.0, style.project_body_size * 1.35)))


def column_groups(shapes: list[Any], tolerance: int = 180000) -> list[list[Any]]:
    groups: list[list[Any]] = []
    for shape in sorted(shapes, key=lambda item: (item.left, item.top)):
        for group in groups:
            if abs(shape.left - group[0].left) <= tolerance:
                group.append(shape)
                break
        else:
            groups.append([shape])
    return [sorted(group, key=lambda item: item.top) for group in sorted(groups, key=lambda group: min(shape.left for shape in group))]


def compact_filled_shape_columns(
    prs: Presentation,
    shapes: list[Any],
    style: TextStyle,
    *,
    gap: int = 90000,
    bottom_margin: int = 260000,
) -> None:
    filled_shapes = [shape for shape in shapes if shape_has_text(shape)]
    if not filled_shapes:
        return

    bottom = prs.slide_height - bottom_margin
    for group in column_groups(filled_shapes):
        y = min(shape.top for shape in group)
        for shape in group:
            estimated_height = estimate_filled_shape_height(shape, style)
            if estimated_height <= 0:
                continue
            min_height = min(max(pt_to_emu(style.body_size * 2.4), 135000), shape.height)
            height = max(min_height, estimated_height + 35000)
            if y + height > bottom:
                height = max(min_height, bottom - y)
            set_shape_geometry(shape, shape.left, y, shape.width, height)
            y += height + gap


def add_vertical_divider(slide: Any, x: int, top: int, bottom: int) -> None:
    if bottom <= top:
        return
    divider = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, top, x, bottom)
    divider.line.width = 12700
    divider.line.color.rgb = rgb("ACACAC")

    # Match the template connector's rounded line caps.
    line = divider._element.spPr.ln
    line.set("cap", "rnd")

    # New connectors inherit the theme's effect style, which adds a drop
    # shadow in PowerPoint and Google Slides. Keep the divider completely flat.
    style = divider._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        divider._element.remove(style)
    divider._element.spPr.append(OxmlElement("a:effectLst"))


def add_horizontal_rule(slide: Any, left: int, top: int, width: int, color: str = AESG_TEAL, weight_pt: float = 1.5) -> Any:
    """Draw a flat horizontal rule on the slide at the requested position."""
    rule = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    rule.line.width = Pt(weight_pt)
    rule.line.color.rgb = rgb(color)
    # Remove the theme effect style so the line does not pick up a shadow.
    style_el = rule.element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style_el is not None:
        rule.element.remove(style_el)
    rule.element.spPr.append(OxmlElement("a:effectLst"))
    return rule


def row_groups(shapes: list[Any], tolerance: int = 120000) -> list[list[Any]]:
    rows: list[list[Any]] = []
    for shape in sorted(shapes, key=lambda item: (item.top, item.left)):
        for row in rows:
            if abs(shape.top - row[0].top) <= tolerance:
                row.append(shape)
                break
        else:
            rows.append([shape])
    return [sorted(row, key=lambda item: item.left) for row in sorted(rows, key=lambda row: min(shape.top for shape in row))]


def fill_project_shape_objects_flowing(
    prs: Presentation,
    slide: Any,
    shapes: list[Any],
    chunks: list[ProjectChunk],
    style: TextStyle,
    *,
    allow_oversize_first: bool = True,
    allow_project_splitting: bool = False,
    height_safety_factor: float = 1.0,
) -> tuple[int, list[Any]]:
    """Fill one editable textbox per project column.

    Project titles, roles, descriptions, and their exact one-line separators
    live in a single text flow. This removes all dependence on the estimated
    bottom of one project box when positioning the next project.
    """
    if not shapes or not chunks:
        for shape in shapes:
            blank_shape(shape)
        return 0, []

    bottom = prs.slide_height - 260000
    columns = column_groups(shapes)
    cursor = 0
    flowing_shapes: list[Any] = []

    if style is LANDSCAPE_STYLE and len(columns) > 1:
        # Populate landscape columns in visual reading order: first project in
        # the left column, second in the right, third back on the left, etc.
        # Each column remains one native editable text flow, so project gaps
        # are still controlled by exact paragraph spacing rather than box
        # coordinates. If the preferred column is full, try the next column
        # before carrying the ordered remainder to a continuation slide.
        prototypes: list[Any] = []
        tops: list[int] = []
        capacities: list[int] = []
        assignments: list[list[ProjectChunk]] = [[] for _column in columns]
        measured_heights = [0 for _column in columns]

        for column in columns:
            for shape in column:
                blank_shape(shape)
            prototype = column[0]
            top = max(min(shape.top for shape in column), int(1.25 * EMU_PER_INCH))
            column_bottom = min(
                bottom,
                max(shape.top + shape.height for shape in column),
            )
            prototypes.append(prototype)
            tops.append(top)
            capacities.append(max(100000, column_bottom - top))

        next_column = 0
        while cursor < len(chunks):
            placed = False
            empty_fallback: tuple[int, int] | None = None
            for offset in range(len(columns)):
                column_index = (next_column + offset) % len(columns)
                candidate = assignments[column_index] + [chunks[cursor]]
                add_projects_flowing_text(prototypes[column_index], candidate, style)
                candidate_height = estimate_project_filled_shape_height(
                    prototypes[column_index],
                    style,
                )
                if (
                    candidate_height * height_safety_factor
                    <= capacities[column_index]
                ):
                    assignments[column_index] = candidate
                    measured_heights[column_index] = candidate_height
                    cursor += 1
                    next_column = (column_index + 1) % len(columns)
                    placed = True
                    break
                if not assignments[column_index] and empty_fallback is None:
                    empty_fallback = (column_index, candidate_height)

            if placed:
                continue
            if (
                empty_fallback is not None
                and allow_project_splitting
                and style in (PORTRAIT_STYLE, LANDSCAPE_STYLE)
            ):
                # The available space can shrink after Experience is rendered
                # and the project lane is moved down. Re-split only this
                # otherwise-oversized project against the final available
                # height, then retry; never force text past the page bottom.
                column_index, _candidate_height = empty_fallback
                replacement = split_project_chunk_to_height(
                    chunks[cursor],
                    prototypes[column_index],
                    capacities[column_index],
                    style,
                )
                if len(replacement) > 1:
                    chunks[cursor : cursor + 1] = replacement
                    continue
            if empty_fallback is not None and allow_oversize_first:
                # The chunking stage normally splits projects taller than a
                # full column. Retain a forward-progress fallback for an
                # indivisible source item without changing project order.
                column_index, candidate_height = empty_fallback
                assignments[column_index] = [chunks[cursor]]
                measured_heights[column_index] = min(
                    capacities[column_index],
                    candidate_height,
                )
                cursor += 1
                next_column = (column_index + 1) % len(columns)
                continue
            break

        for column_index, selected in enumerate(assignments):
            if not selected:
                blank_shape(prototypes[column_index])
                continue
            prototype = prototypes[column_index]
            add_projects_flowing_text(prototype, selected, style)
            set_shape_geometry(
                prototype,
                prototype.left,
                tops[column_index],
                prototype.width,
                max(
                    pt_to_emu(style.project_body_size * 1.35),
                    measured_heights[column_index],
                ),
            )
            flowing_shapes.append(prototype)

        return cursor, flowing_shapes

    for column in columns:
        for shape in column:
            blank_shape(shape)
        if cursor >= len(chunks):
            continue

        prototype = column[0]
        top = min(shape.top for shape in column)
        if style is LANDSCAPE_STYLE:
            # The landscape continuation banner is part of the slide artwork,
            # not a detectable shape. Reserve enough space for the heading
            # below it before measuring how many projects fit in this column.
            # Moving the flow itself avoids fixing the heading by overlapping
            # the first project or silently extending content past the bottom.
            top = max(top, int(1.25 * EMU_PER_INCH))
        left = prototype.left
        width = prototype.width
        # Respect the actual vertical extent of this column. Mixed-layout
        # slides can expose a deliberately short, single-column project lane;
        # using the global slide bottom here made an oversized project appear
        # to fit and then grow beyond the reserved region.
        column_bottom = min(
            bottom,
            max(shape.top + shape.height for shape in column),
        )
        capacity = max(100000, column_bottom - top)
        selected: list[ProjectChunk] = []
        measured_height = 0

        while cursor + len(selected) < len(chunks):
            candidate = selected + [chunks[cursor + len(selected)]]
            add_projects_flowing_text(prototype, candidate, style)
            candidate_height = estimate_project_filled_shape_height(prototype, style)
            if candidate_height * height_safety_factor <= capacity:
                selected = candidate
                measured_height = candidate_height
                continue
            break

        if not selected:
            if allow_project_splitting and style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
                replacement = split_project_chunk_to_height(
                    chunks[cursor], prototype, capacity, style
                )
                if len(replacement) > 1:
                    chunks[cursor : cursor + 1] = replacement
                    # Re-enter the same column with a smaller, independently
                    # measurable fragment rather than accepting a clipped one.
                    return fill_project_shape_objects_flowing(
                        prs,
                        slide,
                        shapes,
                        chunks,
                        style,
                        allow_oversize_first=allow_oversize_first,
                        allow_project_splitting=allow_project_splitting,
                        height_safety_factor=height_safety_factor,
                    )
            if not allow_oversize_first:
                # A short mixed/same-slide region must leave the complete
                # project for the next full-height page. Previously this
                # single-column branch forced the text into the region even
                # when the caller explicitly disabled that fallback.
                blank_shape(prototype)
                continue
            # Oversize projects should already have been sentence-split by the
            # chunking stage. Keep a safe fallback so the renderer can make
            # forward progress if a source contains an indivisible long item.
            selected = [chunks[cursor]]
            add_projects_flowing_text(prototype, selected, style)
            measured_height = min(
                capacity,
                estimate_project_filled_shape_height(prototype, style),
            )
        else:
            add_projects_flowing_text(prototype, selected, style)

        set_shape_geometry(
            prototype,
            left,
            top,
            width,
            max(pt_to_emu(style.project_body_size * 1.35), measured_height),
        )
        flowing_shapes.append(prototype)
        cursor += len(selected)

    return cursor, flowing_shapes


def fill_project_shape_objects_compact(
    prs: Presentation,
    slide: Any,
    shapes: list[Any],
    chunks: list[ProjectChunk],
    style: TextStyle,
) -> int:
    if not shapes:
        return 0

    geometry = {
        id(shape): (shape.left, shape.top, shape.width, shape.height)
        for shape in shapes
    }
    rows = row_groups(shapes)
    cursor = 0
    y = min(geometry[id(shape)][1] for shape in shapes)
    bottom = prs.slide_height - 260000
    gap = compact_project_gap(style)

    for row_index, row in enumerate(rows):
        if cursor >= len(chunks):
            for shape in row:
                blank_shape(shape)
            continue

        row_chunks = chunks[cursor : cursor + len(row)]
        for index, shape in enumerate(row):
            if index < len(row_chunks):
                add_project_text(shape, row_chunks[index], style)
            else:
                blank_shape(shape)
        # Measure the actual formatted paragraphs. Placeholder auto-fit rules
        # can make the pre-render character estimate too short.
        row_heights = [estimate_project_filled_shape_height(shape, style) for shape in row[: len(row_chunks)]]
        original_height = max(geometry[id(shape)][3] for shape in row)
        min_height = compact_project_min_height(style, original_height)
        measurement_padding = 0 if style is PORTRAIT_STYLE else pt_to_emu(1.5)
        row_height = max(min_height, max(row_heights, default=min_height) + measurement_padding)

        if cursor > 0 and y + row_height > bottom:
            for remaining_row in rows[row_index:]:
                for shape in remaining_row:
                    blank_shape(shape)
            return cursor

        if y + row_height > bottom:
            row_height = max(min_height, bottom - y)

        for index, shape in enumerate(row):
            original_left, _original_top, original_width, _original_height = geometry[id(shape)]
            set_shape_geometry(shape, original_left, y, original_width, row_height)

        cursor += len(row_chunks)
        y += row_height + gap

    if cursor < len(chunks) and rows:
        template_row = rows[-1]
        while cursor < len(chunks):
            row_chunks = chunks[cursor : cursor + len(template_row)]
            new_shapes: list[Any] = []
            for shape_index, chunk in enumerate(row_chunks):
                template_shape = template_row[shape_index]
                original_left, _original_top, original_width, _original_height = geometry[id(template_shape)]
                new_shape = slide.shapes.add_textbox(original_left, y, original_width, max(100000, bottom - y))
                add_project_text(new_shape, chunk, style)
                new_shapes.append(new_shape)
            row_heights = [estimate_project_filled_shape_height(shape, style) for shape in new_shapes]
            original_height = max(geometry[id(shape)][3] for shape in template_row)
            min_height = compact_project_min_height(style, original_height)
            measurement_padding = 0 if style is PORTRAIT_STYLE else pt_to_emu(1.5)
            row_height = max(min_height, max(row_heights, default=min_height) + measurement_padding)
            if y + row_height > bottom:
                for shape in new_shapes:
                    parent = shape.element.getparent()
                    if parent is not None:
                        parent.remove(shape.element)
                break

            for new_shape in new_shapes:
                set_shape_geometry(new_shape, new_shape.left, y, new_shape.width, row_height)

            cursor += len(row_chunks)
            y += row_height + gap

    return cursor


def fill_project_shape_objects_masonry(
    prs: Presentation,
    slide: Any,
    shapes: list[Any],
    chunks: list[ProjectChunk],
    style: TextStyle,
    *,
    balance_columns: bool = False,
    allow_project_splitting: bool = False,
) -> tuple[int, list[Any]]:
    """Pack the largest consecutive project prefix across available columns."""
    if not shapes:
        return 0, []

    columns = column_groups(shapes)
    prototypes = [column[0] for column in columns]
    column_top = [min(shape.top for shape in column) for column in columns]
    bottom = prs.slide_height - 260000
    box_padding = pt_to_emu(1.5)
    gap = compact_project_gap(style)
    for shape in shapes:
        blank_shape(shape)

    capacities = [bottom - top for top in column_top]
    max_capacity = max(capacities, default=0)
    min_heights = [compact_project_min_height(style, prototype.height) for prototype in prototypes]

    # Recombine fragments belonging to the same project whenever the complete
    # project can fit in a continuation column. This keeps projects atomic and
    # prevents numbered fragments from being distributed across columns.
    atomic: list[tuple[ProjectChunk, int, int]] = []
    cursor = 0
    while cursor < len(chunks):
        chunk = chunks[cursor]
        group_end = cursor + 1
        identity = project_identity(chunk)
        while group_end < len(chunks) and project_identity(chunks[group_end]) == identity:
            group_end += 1
        group = chunks[cursor:group_end]
        combined = chunk
        if len(group) > 1:
            combined = ProjectChunk(
                project=group[0].project,
                body=clean_text(" ".join(part.body for part in group if part.body)),
                part=1,
                parts=1,
            )

        combined_height = max(
            min_heights[0],
            estimate_compact_project_chunk_height(combined, prototypes[0], style) + box_padding,
        )
        if len(group) > 1 and combined_height <= max_capacity:
            candidates = [(combined, len(group))]
        elif len(group) > 1 and allow_project_splitting:
            # Only a project intrinsically taller than the entire
            # project-only region is allowed to retain its sentence chunks.
            candidates = [(part, 1) for part in group]
        else:
            # In a narrow mixed layout, leave the oversize atomic project for
            # the following full-width project-only slide.
            candidates = [(combined, len(group))] if len(group) > 1 else [(chunk, 1)]

        for placement_chunk, consume_count in candidates:
            height = max(
                min_heights[0],
                estimate_compact_project_chunk_height(placement_chunk, prototypes[0], style) + box_padding,
            )
            atomic.append((placement_chunk, consume_count, height))
        cursor = group_end

    if not atomic or max_capacity <= 0:
        return 0, []

    # Two-column prefix packing. Each state represents the used height in each
    # column after placing every project in the current prefix. Quantizing the
    # heights keeps the state space small while rounding upward ensures the
    # selected layout remains conservative.
    unit = 50000
    capacity_units = [max(0, int((capacity + gap) // unit)) for capacity in capacities]
    states: dict[tuple[int, ...], tuple[int, ...]] = {tuple(0 for _ in columns): tuple()}
    selected_count = 0
    selected_assignment: tuple[int, ...] = tuple()
    for atomic_index, (_chunk, _consume_count, height) in enumerate(atomic):
        weight = max(1, int(math.ceil((height + gap) / unit)))
        next_states: dict[tuple[int, ...], tuple[int, ...]] = {}
        for used_units, assignment in states.items():
            # Preserve column-major reading order by allowing one forward-only
            # split between columns. Balanced continuation pages choose the
            # best split point instead of interleaving projects across columns.
            first_column = assignment[-1] if assignment else 0
            for column_index in range(first_column, len(columns)):
                if used_units[column_index] + weight > capacity_units[column_index]:
                    continue
                updated = list(used_units)
                updated[column_index] += weight
                key = tuple(updated)
                candidate = assignment + (column_index,)
                existing = next_states.get(key)
                if existing is None or candidate < existing:
                    next_states[key] = candidate
        if not next_states:
            break
        states = next_states
        selected_count = atomic_index + 1
        if balance_columns and len(columns) > 1:
            # Maximise use of the emptier column, then minimise imbalance. The
            # number of selected projects is already fixed by the outer loop.
            best_used = max(
                states,
                key=lambda used: (
                    min(used),
                    -abs(max(used) - min(used)),
                    sum(used),
                ),
            )
        else:
            # Prefer fuller earlier columns for front-page reading order.
            best_used = max(states, key=lambda used: tuple(used))
        selected_assignment = states[best_used]

    if selected_count <= 0:
        return 0, []

    created: list[Any] = []
    column_y = list(column_top)
    consumed_chunks = 0
    for atomic_index in range(selected_count):
        column_index = selected_assignment[atomic_index]
        prototype = prototypes[column_index]
        placement_chunk, consume_count, planned_height = atomic[atomic_index]
        available = bottom - column_y[column_index]
        if planned_height > available:
            break
        new_shape = slide.shapes.add_textbox(
            prototype.left,
            column_y[column_index],
            prototype.width,
            available,
        )
        add_project_text(new_shape, placement_chunk, style)
        actual_height = max(
            min_heights[column_index],
            estimate_compact_project_filled_height(new_shape, style) + box_padding,
        )
        # The pre-render estimate is intentionally conservative for packing.
        # Once Office-formatted paragraphs exist, use their measured height so
        # the next project starts about one text line below this one instead of
        # retaining invisible placeholder-era padding.
        height = actual_height
        if height > available:
            parent = new_shape.element.getparent()
            if parent is not None:
                parent.remove(new_shape.element)
            break
        set_shape_geometry(new_shape, new_shape.left, new_shape.top, new_shape.width, height)
        created.append(new_shape)
        column_y[column_index] += height + gap
        consumed_chunks += consume_count
    return consumed_chunks, created


def fill_project_shapes_compact(
    prs: Presentation,
    slide: Any,
    shape_indexes: list[int],
    chunks: list[ProjectChunk],
    style: TextStyle,
) -> int:
    return fill_project_shape_objects_compact(prs, slide, find_textbox_shapes(slide, shape_indexes), chunks, style)


def fill_front_project_shapes(
    prs: Presentation,
    slide: Any,
    shape_indexes: list[int],
    chunks: list[ProjectChunk],
    style: TextStyle,
) -> int:
    shapes = find_textbox_shapes(slide, shape_indexes)
    if not shapes:
        return 0
    max_height = min(shape.height for shape in shapes)
    groups: list[list[ProjectChunk]] = []
    for chunk in chunks:
        if groups and project_identity(groups[-1][0]) == project_identity(chunk):
            groups[-1].append(chunk)
        else:
            groups.append([chunk])

    front_chunks: list[ProjectChunk] = []
    remainder: list[ProjectChunk] = []
    for group in groups:
        if len(front_chunks) >= len(shapes):
            remainder.extend(group)
            continue

        combined = ProjectChunk(
            project=group[0].project,
            body=clean_text(" ".join(chunk.body for chunk in group if chunk.body)),
            part=1,
            parts=1,
        )
        if len(group) == 1 or estimate_project_chunk_height(combined, shapes[0], style) <= max_height:
            front_chunks.append(combined)
        else:
            front_chunks.append(group[0])
            remainder.extend(group[1:])

    used, _ = fill_project_shapes(shapes, front_chunks, style)
    for shape in shapes[used:]:
        cover_static_project_heading(slide, shape)
    chunks[:] = front_chunks + remainder
    return used


def reorder_slides(prs: Presentation, indexes: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    ordered = [slide_ids[index] for index in indexes if index < len(slide_ids)]
    retained_relationship_ids = {slide_id.rId for slide_id in ordered}
    for slide_id in slide_ids:
        slide_id_list.remove(slide_id)
    for slide_id in ordered:
        slide_id_list.append(slide_id)
    # Removing an entry from ``p:sldIdLst`` does not remove its presentation
    # relationship.  If a later duplicate slide reuses that orphan's part
    # name, the deck appears to open normally but becomes an invalid ZIP as
    # soon as another tool edits and saves it. Drop every unused slide
    # relationship so each visible slide has exactly one package part.
    for relationship_id, relationship in list(prs.part.rels.items()):
        if (
            relationship.reltype.endswith("/slide")
            and relationship_id not in retained_relationship_ids
        ):
            prs.part.drop_rel(relationship_id)


def fill_named_section_shape(
    shape: Any,
    heading: str,
    items: list[str],
    style: TextStyle,
) -> list[str]:
    packed_sections, remaining_sections = pack_sections_into_shape([(heading, items, False)], shape, style)
    add_section_text(shape, packed_sections, style) if packed_sections else blank_shape(shape)
    return remaining_sections[0][1] if remaining_sections else []


def fill_standard_sections(slide: Any, cv: dict[str, Any], mapping: dict[str, int], style: TextStyle) -> list[tuple[str, list[str], bool]]:
    sections = standard_remainders(cv)
    original_lengths = {key: len(items) for key, items in sections.items()}

    key_index = mapping.get("key")
    if key_index is not None and key_index < len(slide.shapes):
        key_items = sections.get("key", [])
        if key_items:
            key_shape = slide.shapes[key_index]
            sections["key"] = fill_named_section_shape(key_shape, "Key Expertise", key_items, style)
        else:
            blank_shape(slide.shapes[key_index])

    experience_index = mapping.get("experience")
    if experience_index is not None and experience_index < len(slide.shapes):
        experience_items = sections.get("experience", [])
        if experience_items:
            experience_shape = slide.shapes[experience_index]
            sections["experience"] = fill_named_section_shape(experience_shape, "Experience", experience_items, style)
        else:
            blank_shape(slide.shapes[experience_index])

    qualifications_index = mapping.get("qualifications")
    if qualifications_index is not None and qualifications_index < len(slide.shapes):
        blocks = qualification_blocks(sections)
        if blocks:
            qualifications_shape = slide.shapes[qualifications_index]
            packed_blocks, remaining_blocks = pack_blocks_into_shape(blocks, qualifications_shape, style)
            add_block_items_text(qualifications_shape, packed_blocks, style) if packed_blocks else blank_shape(qualifications_shape)
            apply_qualification_blocks(sections, remaining_blocks)
        else:
            blank_shape(slide.shapes[qualifications_index])
    continuation_sections: list[tuple[str, list[str], bool]] = []
    for key, heading in (
        ("experience", "Experience"),
        ("key", "Key Expertise"),
        ("qualifications", "Qualifications"),
        ("education", "Education"),
        ("memberships", "Memberships"),
    ):
        items = sections.get(key, [])
        if not items:
            continue
        started_on_front = len(items) < original_lengths.get(key, len(items))
        continuation_sections.append((heading, items, started_on_front))
    return continuation_sections


def fill_flowing_side_sections(
    prs: Presentation,
    slide: Any,
    cv: dict[str, Any],
    mapping: dict[str, int],
    style: TextStyle,
    *,
    bottom_margin: int = 260000,
    bottom_buffer: int = 0,
    expand_experience: bool = False,
    density_factor: float = 98.0,
) -> list[tuple[str, list[str], bool]]:
    """Fill the portrait side rail as one flowing region.

    The source templates expose separate expertise and qualification boxes,
    but their fixed heights leave usable page space stranded. Reusing the
    expertise box as a full-height flow region lets Word-length CV sections
    consume that space before a continuation slide is created.
    """
    sections = standard_remainders(cv)
    original_lengths = {key: len(items) for key, items in sections.items()}

    key_index = mapping.get("key")
    qualifications_index = mapping.get("qualifications")
    if key_index is not None and key_index < len(slide.shapes):
        flow_shape = slide.shapes[key_index]
        flow_bottom = max(flow_shape.top + 100000, prs.slide_height - bottom_margin)
        set_shape_geometry(
            flow_shape,
            flow_shape.left,
            flow_shape.top,
            flow_shape.width,
            max(100000, flow_bottom - flow_shape.top - bottom_buffer),
        )
        flowing_sections = [
            (heading, sections.get(key, []), False)
            for key, heading in (
                ("key", "Key Expertise"),
                ("qualifications", "Qualifications"),
                ("education", "Education"),
                ("memberships", "Memberships"),
            )
            if sections.get(key)
        ]
        packed, remaining = pack_sections_into_shape(
            flowing_sections,
            flow_shape,
            style,
            density_factor=density_factor,
            keep_sections_together=style is not PORTRAIT_STYLE,
        )
        add_section_text(flow_shape, packed, style) if packed else blank_shape(flow_shape)

        for key in ("key", "qualifications", "education", "memberships"):
            sections[key] = []
        heading_to_key = {
            "Key Expertise": "key",
            "Qualifications": "qualifications",
            "Education": "education",
            "Memberships": "memberships",
        }
        for heading, items, _continued in remaining:
            sections[heading_to_key[heading]] = items

    if qualifications_index is not None and qualifications_index < len(slide.shapes):
        blank_shape(slide.shapes[qualifications_index])

    experience_index = mapping.get("experience")
    if experience_index is not None and experience_index < len(slide.shapes):
        experience_shape = slide.shapes[experience_index]
        if expand_experience:
            experience_bottom = max(experience_shape.top + 100000, prs.slide_height - bottom_margin)
            set_shape_geometry(
                experience_shape,
                experience_shape.left,
                experience_shape.top,
                experience_shape.width,
                max(100000, experience_bottom - experience_shape.top - bottom_buffer),
            )
        experience_items = sections.get("experience", [])
        if experience_items:
            sections["experience"] = fill_named_section_shape(
                experience_shape,
                "Experience",
                experience_items,
                style,
            )
        else:
            blank_shape(slide.shapes[experience_index])

    continuation_sections: list[tuple[str, list[str], bool]] = []
    for key, heading in (
        ("experience", "Experience"),
        ("key", "Key Expertise"),
        ("qualifications", "Qualifications"),
        ("education", "Education"),
        ("memberships", "Memberships"),
    ):
        items = sections.get(key, [])
        if not items:
            continue
        started_on_front = len(items) < original_lengths.get(key, len(items))
        continuation_sections.append((heading, items, started_on_front))
    return continuation_sections


def fill_flowing_all_sections(
    prs: Presentation,
    slide: Any,
    cv: dict[str, Any],
    mapping: dict[str, int],
    style: TextStyle,
    *,
    bottom_margin: int = 260000,
    bottom_buffer: int = 0,
    density_factor: float = 98.0,
) -> list[tuple[str, list[str], bool]]:
    """Flow every landscape side section through one ordered column."""
    sections = standard_remainders(cv)
    original_lengths = {key: len(items) for key, items in sections.items()}
    valid_indexes = [index for index in mapping.values() if index is not None and index < len(slide.shapes)]
    if not valid_indexes:
        return section_data(cv)

    mapped_shapes = [slide.shapes[index] for index in valid_indexes]
    anchor_index = mapping.get("key", valid_indexes[0])
    if anchor_index is None or anchor_index >= len(slide.shapes):
        anchor_index = valid_indexes[0]
    flow_shape = slide.shapes[anchor_index]
    flow_top = min(shape.top for shape in mapped_shapes)
    flow_left = min(shape.left for shape in mapped_shapes)
    flow_right = max(shape.left + shape.width for shape in mapped_shapes)
    flow_bottom = max(flow_top + 100000, prs.slide_height - bottom_margin)
    set_shape_geometry(
        flow_shape,
        flow_left,
        flow_top,
        flow_right - flow_left,
        max(100000, flow_bottom - flow_top - bottom_buffer),
    )

    flowing_sections = [
        (heading, sections.get(key, []), False)
        for key, heading in (
            ("experience", "Experience"),
            ("key", "Key Expertise"),
            ("qualifications", "Qualifications"),
            ("education", "Education"),
            ("memberships", "Memberships"),
        )
        if sections.get(key)
    ]
    packed, remaining = pack_sections_into_shape(
        flowing_sections,
        flow_shape,
        style,
        density_factor=density_factor,
    )
    add_section_text(flow_shape, packed, style) if packed else blank_shape(flow_shape)

    for index in valid_indexes:
        if index != anchor_index:
            blank_shape(slide.shapes[index])

    heading_to_key = {
        "Experience": "experience",
        "Key Expertise": "key",
        "Qualifications": "qualifications",
        "Education": "education",
        "Memberships": "memberships",
    }
    for key in sections:
        sections[key] = []
    for heading, items, _continued in remaining:
        sections[heading_to_key[heading]] = items

    continuation_sections: list[tuple[str, list[str], bool]] = []
    for key, heading in (
        ("experience", "Experience"),
        ("key", "Key Expertise"),
        ("qualifications", "Qualifications"),
        ("education", "Education"),
        ("memberships", "Memberships"),
    ):
        items = sections.get(key, [])
        if not items:
            continue
        started_on_front = len(items) < original_lengths.get(key, len(items))
        continuation_sections.append((heading, items, started_on_front))
    return continuation_sections


def estimate_textbox_height_emu(text: str, shape: Any, size: float, *, factor: float = 108.0, leading: float = 1.16) -> int:
    return pt_to_emu(paragraph_height(text, shape.width, size, factor=factor, leading=leading) + 8)


def compact_below_overview(
    slide: Any,
    overview_index: int,
    content_indexes: Iterable[int],
    style: TextStyle,
    *,
    margin: int = 240000,
    minimum_content_top: int | None = None,
) -> None:
    if overview_index >= len(slide.shapes):
        return
    overview_shape = slide.shapes[overview_index]
    content_shapes = [slide.shapes[index] for index in content_indexes if index < len(slide.shapes)]
    if not content_shapes:
        return
    if not clean_text(overview_shape.text):
        return
    overview_height_pt = 0.0
    for paragraph in overview_shape.text_frame.paragraphs:
        text = paragraph_text(paragraph)
        if not text:
            continue
        size = paragraph_font_size(paragraph, style.overview_size)
        if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
            overview_height_pt += verdana_paragraph_height(
                text,
                overview_shape.width,
                size,
                leading=1.14,
            )
        else:
            overview_height_pt += paragraph_height(
                text,
                overview_shape.width,
                size,
                factor=128.0,
                leading=1.14,
            )
        if paragraph.space_after:
            overview_height_pt += paragraph.space_after.pt
    overview_height = pt_to_emu(overview_height_pt + 8)
    # Portrait V1 deliberately uses a fixed 9pt Verdana body.  Do not allow a
    # long overview on the alternate front layout to be made legible only by
    # shrinking it: grow its editable frame, then make room beneath it.  The
    # following section/project packing will naturally move excess projects
    # to a continuation slide.
    if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE) and overview_height > overview_shape.height:
        set_shape_geometry(
            overview_shape,
            overview_shape.left,
            overview_shape.top,
            overview_shape.width,
            overview_height,
        )
    desired_top = overview_shape.top + overview_height + margin
    if minimum_content_top is not None:
        desired_top = max(desired_top, minimum_content_top)
    # The second landscape front layout keeps its full-width divider on the
    # slide layout rather than as an ordinary slide shape.  Because layout
    # shapes are not present in ``slide.shapes``, overview compaction used to
    # pull the section boxes through that fixed rule.  Honour the template
    # rule as an immutable lower boundary and retain the template's optical
    # heading clearance beneath it.
    if style is LANDSCAPE_STYLE:
        template_rule_shapes: list[Any] = []
        layout = getattr(slide, "slide_layout", None)
        if layout is not None:
            template_rule_shapes.extend(layout.shapes)
            master = getattr(layout, "slide_master", None)
            if master is not None:
                template_rule_shapes.extend(master.shapes)
        template_horizontal_rules = [
            shape
            for shape in template_rule_shapes
            if getattr(shape, "shape_type", None) == 9
            and abs(shape.height) < 30000
            and shape.width > int(2.0 * EMU_PER_INCH)
            and shape.top > overview_shape.top
        ]
        if template_horizontal_rules:
            rule = max(template_horizontal_rules, key=lambda shape: shape.width)
            # Remove the fixed-position layout rule and redraw a slide-level rule
            # so the overview separator is always gray and sits below the actual
            # rendered overview instead of being covered by long text.
            parent = rule.element.getparent()
            if parent is not None:
                parent.remove(rule.element)
            add_horizontal_rule(
                slide,
                rule.left,
                max(rule.top, overview_shape.top + overview_shape.height + 155000),
                rule.width,
                color="ACACAC",
                weight_pt=1.0,
            )
    # Some alternatives have a full-width horizontal rule below the profile
    # area. Keep section headings optically clear of that rule after compaction.
    horizontal_rules = [
        shape
        for shape in slide.shapes
        if getattr(shape, "shape_type", None) == 9
        and abs(shape.height) < 30000
        and shape.width > int(2.0 * EMU_PER_INCH)
        and shape.top > overview_shape.top
    ]
    if horizontal_rules:
        natural_rule_top = overview_shape.top + overview_shape.height + 155000
        if minimum_content_top is not None:
            natural_rule_top = max(natural_rule_top, minimum_content_top - 145000)
        for rule in horizontal_rules:
            if rule.top > natural_rule_top:
                set_shape_geometry(rule, rule.left, natural_rule_top, rule.width, rule.height)
        # Section text boxes can render a few points above their nominal
        # bounds.  Give the template's horizontal divider a full text-leading
        # clearance so a heading can never visually touch or cross it.
        desired_top = max(desired_top, max(shape.top for shape in horizontal_rules) + 235000)
    current_top = min(shape.top for shape in content_shapes)
    # Position the following content from the overview's rendered height, not
    # from the bottom of its generously sized template placeholder. This keeps
    # the source template's visual gap while allowing short portrait overviews
    # to pull Experience/Projects upward. Long overviews still push everything
    # down because ``desired_top`` is measured from their actual line count.
    shift = desired_top - current_top
    if not shift:
        return
    if style is not PORTRAIT_STYLE:
        shift = min(shift, 1250000)
    for shape in content_shapes:
        left, width, height = shape.left, shape.width, shape.height
        set_shape_geometry(shape, left, max(0, shape.top + shift), width, height)


def render_profile(
    slide: Any,
    cv: dict[str, Any],
    tmpdir: Path,
    style: TextStyle,
    mapping: dict[str, int],
    *,
    name_color: str = WHITE,
    role_color: str = WHITE,
    photo_background: str = AESG_TEAL,
) -> None:
    first, last, role = profile_parts(cv)
    if mapping["photo"] < len(slide.shapes):
        replace_photo(slide.shapes[mapping["photo"]], cv, tmpdir, photo_background)
    if mapping["name"] < len(slide.shapes):
        name_shape = slide.shapes[mapping["name"]]
        if mapping["overview"] < len(slide.shapes):
            overview_shape = slide.shapes[mapping["overview"]]
            available_right = overview_shape.left + overview_shape.width
            if available_right > name_shape.left + name_shape.width:
                set_shape_geometry(
                    name_shape,
                    name_shape.left,
                    name_shape.top,
                    available_right - name_shape.left,
                    name_shape.height,
                )
        fill_name(
            name_shape,
            first,
            last,
            size=style.name_size,
            color=name_color,
            font=style.title_font,
        )
    if mapping["role"] < len(slide.shapes):
        fill_plain(
            slide.shapes[mapping["role"]],
            role,
            size=style.role_size,
            color=role_color,
            font=style.body_font,
        )
    if mapping["overview"] < len(slide.shapes):
        overview_shape = slide.shapes[mapping["overview"]]
        overview = clean_text(cv.get("overview") or cv.get("overview_full") or "")
        overview_base = replace(style, body_size=style.overview_size)
        overview_estimate = paragraph_height(overview, overview_shape.width, style.overview_size, factor=108.0, leading=1.16)
        # The new portrait and landscape templates both use fixed 9pt
        # Verdana body text.  Preserve that typography and grow the editable
        # overview frame when needed; subsequent content packing can then
        # continue naturally instead of silently shrinking the overview.
        overview_style = (
            overview_base
            if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE)
            else fitted_style_for_height(
                overview_shape,
                overview_base,
                overview_estimate,
                min_body=6.2,
                min_heading=6.2,
            )
        )
        fill_plain(
            overview_shape,
            overview,
            size=overview_style.body_size,
            font=style.body_font,
            align=PP_ALIGN.JUSTIFY,
        )
        if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
            rendered_height = estimate_textbox_height_emu(
                overview,
                overview_shape,
                style.overview_size,
            )
            if rendered_height > overview_shape.height:
                set_shape_geometry(
                    overview_shape,
                    overview_shape.left,
                    overview_shape.top,
                    overview_shape.width,
                    rendered_height,
                )

    # Portrait front slides (both main and alternate) should have a small
    # teal underline beneath the role, matching the template design.
    if style is PORTRAIT_STYLE and mapping["role"] < len(slide.shapes) and mapping["overview"] < len(slide.shapes):
        role_shape = slide.shapes[mapping["role"]]
        overview_shape = slide.shapes[mapping["overview"]]
        # Remove any existing short template title line from the layout/masters
        # so we don't end up with two lines when the layout already has one.
        layout = getattr(slide, "slide_layout", None)
        template_line_shapes: list[Any] = []
        if layout is not None:
            template_line_shapes.extend(layout.shapes)
            master = getattr(layout, "slide_master", None)
            if master is not None:
                template_line_shapes.extend(master.shapes)
        for line in template_line_shapes:
            if (
                getattr(line, "shape_type", None) == 9
                and abs(line.height) < 30000
                and line.width < int(2.0 * EMU_PER_INCH)
                and line.top > role_shape.top
                and line.top < overview_shape.top
                and line.left > 1000000
            ):
                parent = line.element.getparent()
                if parent is not None:
                    parent.remove(line.element)
        # Draw the title line centred between the role and the overview.
        line_top = (role_shape.top + role_shape.height + overview_shape.top) // 2
        add_horizontal_rule(
            slide,
            role_shape.left,
            line_top,
            914400,
            color=AESG_TEAL,
            weight_pt=1.5,
        )


def render_front_option(
    prs: Presentation,
    slide: Any,
    cv: dict[str, Any],
    tmpdir: Path,
    style: TextStyle,
    profile_mapping: dict[str, int],
    section_mapping: dict[str, int],
    project_indexes: list[int],
    *,
    name_color: str = WHITE,
    role_color: str = WHITE,
    photo_background: str = AESG_TEAL,
    flowing_side_sections: bool = False,
    flowing_all_sections: bool = False,
    expand_experience: bool = False,
    precompact_before_sections: bool = False,
    content_respects_photo: bool = False,
    compact_projects_with_sections: bool = True,
    section_density_factor: float = 98.0,
    project_masonry: bool = False,
    balance_section_overflow: bool = False,
    align_project_heading_with_sections: bool = False,
    front_divider_x: int | None = None,
    front_divider_bottom: int | None = None,
    section_bottom_margin: int = 260000,
    section_bottom_buffer: int = 0,
) -> tuple[list[tuple[str, list[str], bool]], list[ProjectChunk]]:
    render_profile(
        slide,
        cv,
        tmpdir,
        style,
        profile_mapping,
        name_color=name_color,
        role_color=role_color,
        photo_background=photo_background,
    )
    overview_index = profile_mapping.get("overview")
    minimum_content_top = None
    if content_respects_photo:
        photo_index = profile_mapping.get("photo")
        if photo_index is not None and photo_index < len(slide.shapes):
            photo_shape = slide.shapes[photo_index]
            minimum_content_top = photo_shape.top + photo_shape.height + 320000
    content_indexes = [
        index
        for index in list(section_mapping.values()) + list(project_indexes)
        if index is not None
    ]
    if precompact_before_sections and overview_index is not None:
        compact_below_overview(
            slide,
            overview_index,
            content_indexes,
            style,
            minimum_content_top=minimum_content_top,
        )

    if flowing_all_sections:
        section_overflow = fill_flowing_all_sections(
            prs,
            slide,
            cv,
            section_mapping,
            style,
            bottom_margin=section_bottom_margin,
            bottom_buffer=section_bottom_buffer,
            density_factor=section_density_factor,
        )
    elif flowing_side_sections:
        section_overflow = fill_flowing_side_sections(
            prs,
            slide,
            cv,
            section_mapping,
            style,
            expand_experience=expand_experience,
            bottom_margin=section_bottom_margin,
            bottom_buffer=section_bottom_buffer,
            density_factor=section_density_factor,
        )
    else:
        section_overflow = fill_standard_sections(slide, cv, section_mapping, style)
    project_shapes = find_textbox_shapes(slide, project_indexes)
    project_records = selected_projects(cv)
    if not project_shapes:
        return section_overflow, project_chunks_for_shapes(
            project_records,
            slide.shapes[0],
            style,
        )
    # Front-page cards follow the reference's "brief scope" treatment. Keep a
    # separate full-description list for continuation slides so the number of
    # projects fitted on the front never changes or truncates later content.
    chunks = project_chunks_for_shapes(
        project_records,
        project_shapes[0],
        style,
        use_teaser=True,
    )
    continuation_chunks = project_chunks_for_shapes(
        project_records,
        project_shapes[0],
        style,
    )
    # Seed the template boxes so their measured heights can participate in the
    # initial compaction. The final pass below repacks the original ordered
    # chunks and may add further boxes.
    used, _ = fill_project_shapes(project_shapes, chunks, style)
    if overview_index is not None and not precompact_before_sections:
        compact_below_overview(
            slide,
            overview_index,
            content_indexes,
            style,
            minimum_content_top=minimum_content_top,
        )
    section_shapes = [
        slide.shapes[index]
        for index in section_mapping.values()
        if index is not None and index < len(slide.shapes)
    ]
    content_shapes = list(section_shapes)
    content_shapes.extend(project_shapes)
    if flowing_all_sections or flowing_side_sections:
        # Flow shapes were already sized with a bottom buffer; don't collapse
        # that reserved space. Compact projects only if requested.
        if compact_projects_with_sections:
            compact_filled_shape_columns(
                prs,
                project_shapes,
                style,
                gap=pt_to_emu(max(7.0, style.body_size * 1.05)),
            )
    else:
        compact_shapes = content_shapes if compact_projects_with_sections else section_shapes
        compact_filled_shape_columns(
            prs,
            compact_shapes,
            style,
            gap=pt_to_emu(max(7.0, style.body_size * 1.05)),
        )
    project_heading_top = None
    if align_project_heading_with_sections:
        filled_sections = [shape for shape in section_shapes if shape_has_text(shape)]
        if filled_sections:
            project_heading_top = min(shape.top for shape in filled_sections)
            project_content_top = project_heading_top + pt_to_emu(style.heading_size * 2.7)
            current_project_top = min(shape.top for shape in project_shapes)
            offset = project_content_top - current_project_top
            for shape in project_shapes:
                set_shape_geometry(
                    shape,
                    shape.left,
                    shape.top + offset,
                    shape.width,
                    shape.height,
                )
    if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
        # Some variants place Experience directly above the project lane,
        # while others place the two areas side by side. Only sections that
        # actually share horizontal space with a project column should push
        # the project heading down. This uses the final rendered height rather
        # than the original placeholder geometry, so a wrapped experience row
        # cannot be covered by SELECTED PROJECTS.
        filled_sections = [shape for shape in section_shapes if shape_has_text(shape)]
        shared_lane_sections = []
        for section_shape in filled_sections:
            for project_shape in project_shapes:
                overlap = max(
                    0,
                    min(section_shape.left + section_shape.width, project_shape.left + project_shape.width)
                    - max(section_shape.left, project_shape.left),
                )
                if overlap >= min(section_shape.width, project_shape.width) * 0.35:
                    shared_lane_sections.append(section_shape)
                    break
        if shared_lane_sections:
            section_bottom = max(
                shape.top + estimate_filled_shape_height(shape, style)
                for shape in shared_lane_sections
            )
            shared_heading_top = section_bottom + pt_to_emu(max(4.0, style.body_size))
            desired_project_top = shared_heading_top + pt_to_emu(style.heading_size * 2.7)
            current_project_top = min(shape.top for shape in project_shapes)
            offset = max(0, desired_project_top - current_project_top)
            if offset:
                for shape in project_shapes:
                    set_shape_geometry(
                        shape,
                        shape.left,
                        shape.top + offset,
                        shape.width,
                        shape.height,
                    )
            project_heading_top = max(project_heading_top or 0, shared_heading_top)
    # The template provides only a small fixed number of front-page project
    # placeholders. Once the overview and experience have been compacted, use
    # the actual remaining column height to add further project blocks before
    # creating a continuation slide.
    max_project_height = prs.slide_height - min(shape.top for shape in project_shapes) - 180000
    chunks = merge_project_chunks_for_shape(chunks, project_shapes[0], style, max_project_height)
    display_project_shapes = project_shapes
    created_project_shapes: list[Any] = []
    if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
        compact_used, flowing_project_shapes = fill_project_shape_objects_flowing(
            prs,
            slide,
            project_shapes,
            chunks,
            style,
        )
        if flowing_project_shapes:
            display_project_shapes = flowing_project_shapes
    elif project_masonry:
        compact_used, created_project_shapes = fill_project_shape_objects_masonry(
            prs,
            slide,
            project_shapes,
            chunks,
            style,
        )
        if created_project_shapes:
            display_project_shapes = created_project_shapes
    else:
        compact_used = fill_project_shape_objects_compact(prs, slide, project_shapes, chunks, style)
    if compact_used > 0:
        used = compact_used
    if (
        balance_section_overflow
        and section_overflow
        and created_project_shapes
        and used == len(chunks)
        and used > 2
    ):
        # Avoid a nearly empty continuation slide containing only a short
        # left-column remainder. Keep enough projects for that page to use the
        # right-hand columns while preserving the original project order.
        overflow_items = sum(len(items) for _heading, items, _continued in section_overflow)
        reserve_count = min(used - 1, 2 if overflow_items >= 4 else 1)
        for shape in created_project_shapes[-reserve_count:]:
            parent = shape.element.getparent()
            if parent is not None:
                parent.remove(shape.element)
        created_project_shapes = created_project_shapes[:-reserve_count]
        display_project_shapes = created_project_shapes
        used -= reserve_count
    if front_divider_x is not None:
        filled_content = [shape for shape in content_shapes if shape_has_text(shape)]
        # The final project shapes may differ from the original placeholders
        # (e.g., flowing/masonry shapes created by the renderer). Include them
        # so the divider stops at the actual bottom of the rendered content.
        filled_content.extend(
            shape for shape in display_project_shapes if shape_has_text(shape)
        )
        # Deduplicate by identity in case the original placeholders were reused.
        filled_content = list({id(shape): shape for shape in filled_content}.values())
        if filled_content:
            divider_top = min(shape.top for shape in filled_content)
            actual_bottom = max(
                shape.top + estimate_filled_shape_height(shape, style)
                for shape in filled_content
            )
            divider_bottom = min(
                front_divider_bottom or (prs.slide_height - 260000),
                actual_bottom + int(0.15 * EMU_PER_INCH),
            )
            if divider_bottom > divider_top:
                add_vertical_divider(slide, front_divider_x, divider_top, divider_bottom)
    cover_layout_labels(slide)
    add_project_area_heading(
        slide,
        display_project_shapes,
        style,
        heading_top=project_heading_top,
    )
    return section_overflow, continuation_chunks[used:]


def fill_section_shapes(shapes: list[Any], sections: list[tuple[str, list[str], bool]], style: TextStyle) -> tuple[list[tuple[str, list[str], bool]], list[list[tuple[str, list[str], bool]]]]:
    keep_sections_together = style is not PORTRAIT_STYLE
    packed = pack_sections(
        sections,
        shapes,
        style,
        keep_sections_together=keep_sections_together,
    )
    for shape, chunks in zip(shapes, packed):
        if chunks:
            add_section_text(shape, chunks, style)
        else:
            blank_shape(shape)
    consumed_count = 0
    remaining = sections
    for shape_chunks in packed:
        if not shape_chunks:
            continue
        _packed, remaining = pack_sections_into_shape(
            remaining,
            shapes[consumed_count] if consumed_count < len(shapes) else shapes[-1],
            style,
            keep_sections_together=keep_sections_together,
        )
        consumed_count += 1
    return remaining, packed


def render_front_portrait(prs: Presentation, cv: dict[str, Any], tmpdir: Path) -> tuple[dict[str, list[str]], list[ProjectChunk]]:
    slide = prs.slides[0]
    first, last, role = profile_parts(cv)
    replace_photo(slide.shapes[0], cv, tmpdir)
    fill_name(
        slide.shapes[1], first, last,
        size=PORTRAIT_STYLE.name_size, color=WHITE, font=PORTRAIT_STYLE.body_font,
    )
    fill_plain(
        slide.shapes[3], role,
        size=PORTRAIT_STYLE.role_size, color=WHITE, font=PORTRAIT_STYLE.body_font,
    )
    fill_plain(
        slide.shapes[2],
        cv.get("overview") or cv.get("overview_full") or "",
        size=PORTRAIT_STYLE.overview_size,
        font=PORTRAIT_STYLE.body_font,
        align=PP_ALIGN.JUSTIFY,
    )

    remaining_sections = standard_remainders(cv)
    packed, remaining = pack_items_into_shape(remaining_sections["key"], slide.shapes[4], PORTRAIT_STYLE)
    add_items_text(slide.shapes[4], packed, PORTRAIT_STYLE) if packed else blank_shape(slide.shapes[4])
    remaining_sections["key"] = remaining

    qual_blocks = qualification_blocks(remaining_sections)
    packed_blocks, remaining_blocks = pack_blocks_into_shape(qual_blocks, slide.shapes[5], PORTRAIT_STYLE)
    add_block_items_text(slide.shapes[5], packed_blocks, PORTRAIT_STYLE) if packed_blocks else blank_shape(slide.shapes[5])
    apply_qualification_blocks(remaining_sections, remaining_blocks)

    packed, remaining = pack_items_into_shape(remaining_sections["experience"], slide.shapes[6], PORTRAIT_STYLE)
    add_items_text(slide.shapes[6], packed, PORTRAIT_STYLE) if packed else blank_shape(slide.shapes[6])
    remaining_sections["experience"] = remaining

    projects = selected_projects(cv)
    project_shapes = [slide.shapes[7], slide.shapes[8]]
    prototype = project_shapes[0]
    chunks = project_chunks_for_shapes(projects, prototype, PORTRAIT_STYLE)
    used, _ = fill_project_shapes(project_shapes, chunks, PORTRAIT_STYLE)
    return remaining_sections, chunks[used:]


def render_project_slides(
    prs: Presentation,
    template_index: int,
    chunks: list[ProjectChunk],
    shape_indexes: list[int],
    style: TextStyle,
    *,
    keep_blank: bool = False,
    masonry: bool = False,
) -> list[int]:
    template_slide = prs.slides[template_index]
    template_layout = template_slide.slide_layout
    template_elements = [copy.deepcopy(shape.element) for shape in template_slide.shapes]
    prototype_shapes = find_textbox_shapes(prs.slides[template_index], shape_indexes)
    if chunks and prototype_shapes:
        max_height = prs.slide_height - min(shape.top for shape in prototype_shapes) - 260000
        chunks = merge_project_chunks_for_shape(chunks, prototype_shapes[0], style, max_height)
    if not chunks:
        if keep_blank:
            slide = prs.slides[template_index]
            for shape in find_textbox_shapes(slide, shape_indexes):
                blank_shape(shape)
            return [template_index]
        return []
    slide_indexes: list[int] = []
    cursor = 0
    first = True
    while cursor < len(chunks):
        if first:
            slide = prs.slides[template_index]
            slide_index = template_index
            first = False
        else:
            slide = duplicate_from_snapshot(prs, template_layout, template_elements)
            slide_index = len(prs.slides) - 1
        display_shapes = find_textbox_shapes(slide, shape_indexes)
        if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
            used, flowing_shapes = fill_project_shape_objects_flowing(
                prs,
                slide,
                display_shapes,
                chunks[cursor:],
                style,
                allow_project_splitting=True,
            )
            if flowing_shapes:
                display_shapes = flowing_shapes
        elif masonry:
            used, created_shapes = fill_project_shape_objects_masonry(
                prs,
                slide,
                display_shapes,
                chunks[cursor:],
                style,
                balance_columns=True,
            )
            if created_shapes:
                display_shapes = created_shapes
        else:
            used = fill_project_shapes_compact(prs, slide, shape_indexes, chunks[cursor:], style)
        if used <= 0:
            shapes = find_textbox_shapes(slide, shape_indexes)
            used, _ = fill_project_shapes(shapes, chunks[cursor:], style)
        if used <= 0:
            break
        cover_layout_labels(slide)
        has_banner = slide_layout_has_teal_banner(slide)
        add_project_area_heading(slide, display_shapes, style, with_rule=not has_banner)
        cursor += used
        slide_indexes.append(slide_index)
    return slide_indexes


def render_mixed_continuation_slides(
    prs: Presentation,
    template_index: int,
    sections: list[tuple[str, list[str], bool]],
    chunks: list[ProjectChunk],
    shape_indexes: list[int],
    style: TextStyle,
    *,
    keep_blank: bool = False,
    section_shape_indexes: list[int] | None = None,
    project_shape_indexes: list[int] | None = None,
    split_layout: SplitContinuationLayout | None = None,
) -> list[int]:
    if not sections and not chunks:
        return []
    if not sections and (
        split_layout is None or not split_layout.project_only_full_width
    ):
        return render_project_slides(
            prs,
            template_index,
            chunks,
            shape_indexes,
            style,
            keep_blank=False,
            masonry=bool(split_layout and split_layout.masonry_projects),
        )

    template_slide = prs.slides[template_index]
    template_layout = template_slide.slide_layout
    template_elements = [copy.deepcopy(shape.element) for shape in template_slide.shapes]
    prototype_indexes = project_shape_indexes or shape_indexes
    prototype_shapes = find_textbox_shapes(prs.slides[template_index], prototype_indexes)
    if chunks and prototype_shapes and split_layout is None:
        max_height = prs.slide_height - min(shape.top for shape in prototype_shapes) - 260000
        chunks = merge_project_chunks_for_shape(chunks, prototype_shapes[0], style, max_height)

    slide_indexes: list[int] = []
    remaining_sections = sections
    cursor = 0
    first = True
    while remaining_sections or cursor < len(chunks):
        before_sections = sum(len(items) for _heading, items, _continued in remaining_sections)
        before_cursor = cursor
        if first:
            slide = prs.slides[template_index]
            slide_index = template_index
            first = False
        else:
            slide = duplicate_from_snapshot(prs, template_layout, template_elements)
            slide_index = len(prs.slides) - 1

        # When the sidebar is exhausted, use the entire content area for
        # projects. Portrait keeps one full-width flowing column; landscape
        # retains its configured two-column grid. The duplicated source slide
        # keeps its native artwork and logo, and all boxes remain editable.
        if (
            not remaining_sections
            and split_layout is not None
            and split_layout.project_only_full_width
            and cursor < len(chunks)
        ):
            for shape in find_textbox_shapes(slide, shape_indexes):
                blank_shape(shape)
            columns = max(1, split_layout.project_columns)
            project_left = split_layout.section_left
            total_gap = split_layout.column_gap * (columns - 1)
            column_width = max(
                100000,
                int((split_layout.right - project_left - total_gap) / columns),
            )
            project_shapes: list[Any] = []
            for column in range(columns):
                left = project_left + column * (column_width + split_layout.column_gap)
                project_shapes.append(
                    slide.shapes.add_textbox(
                        left,
                        split_layout.top,
                        column_width,
                        split_layout.bottom - split_layout.top,
                    )
                )
            merged = merge_project_chunks_for_shape(
                chunks[cursor:],
                project_shapes[0],
                style,
                split_layout.bottom - split_layout.top,
            )
            chunks = chunks[:cursor] + merged
            if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
                used, project_shapes = fill_project_shape_objects_flowing(
                    prs,
                    slide,
                    project_shapes,
                    chunks[cursor:],
                    style,
                    allow_project_splitting=True,
                )
            else:
                used, project_shapes = fill_project_shape_objects_masonry(
                    prs,
                    slide,
                    project_shapes,
                    chunks[cursor:],
                    style,
                    balance_columns=True,
                    allow_project_splitting=True,
                )
            cursor += used
            cover_layout_labels(slide)
            has_banner = slide_layout_has_teal_banner(slide)
            add_project_area_heading(
                slide,
                project_shapes,
                style,
                project_only_full_width=True,
                with_rule=not has_banner,
            )
            slide_indexes.append(slide_index)
            if used <= 0:
                break
            continue

        had_sections_on_slide = bool(remaining_sections)
        # A divider represents real left-side content on this slide. Merely
        # configuring a divider must not reserve an empty sidebar after those
        # sections have been fully consumed.
        use_split_layout = split_layout is not None and bool(remaining_sections)
        if use_split_layout:
            # The continuation page has one flowing left column and a separate
            # flexible project grid. This preserves the first slide's divider
            # position and prevents narrow or overlapping project boxes.
            for shape in find_textbox_shapes(slide, shape_indexes):
                blank_shape(shape)

            section_width = max(100000, split_layout.divider_x - split_layout.section_left - 90000)
            section_shape = slide.shapes.add_textbox(
                split_layout.section_left,
                split_layout.top,
                section_width,
                split_layout.bottom - split_layout.top,
            )
            packed, remaining_sections = pack_sections_into_shape(
                remaining_sections,
                section_shape,
                style,
                density_factor=split_layout.section_density_factor,
                keep_sections_together=style is not PORTRAIT_STYLE,
            )
            add_section_text(section_shape, packed, style) if packed else blank_shape(section_shape)

            # A landscape continuation slide may change layout after the last
            # sidebar block. Cap the narrow project lane at the rendered
            # sidebar bottom, then use the remaining lower region as a native
            # full-width two-column project area. If sidebar content still
            # remains for another slide, retain the split layout to the bottom.
            narrow_bottom = split_layout.bottom
            allow_same_slide_full_width = (
                style is LANDSCAPE_STYLE
                and split_layout.project_only_full_width
                and not remaining_sections
                and bool(packed)
            )
            if allow_same_slide_full_width:
                rendered_section_height = estimate_filled_shape_height(section_shape, style)
                narrow_bottom = min(
                    split_layout.bottom,
                    section_shape.top
                    + max(
                        rendered_section_height,
                        pt_to_emu(style.heading_size * 2.5),
                    )
                    + int(0.12 * EMU_PER_INCH),
                )

            project_shapes: list[Any] = []
            columns = max(1, split_layout.project_columns)
            total_gap = split_layout.column_gap * (columns - 1)
            column_width = max(
                100000,
                int((split_layout.right - split_layout.project_left - total_gap) / columns),
            )
            for column in range(columns):
                left = split_layout.project_left + column * (column_width + split_layout.column_gap)
                project_shapes.append(
                    slide.shapes.add_textbox(
                        left,
                        split_layout.top,
                        column_width,
                        max(100000, narrow_bottom - split_layout.top),
                    )
                )

            if cursor < len(chunks) and project_shapes:
                merged = merge_project_chunks_for_shape(
                    chunks[cursor:],
                    project_shapes[0],
                    style,
                    split_layout.bottom - split_layout.top,
                )
                chunks = chunks[:cursor] + merged
            if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
                used, project_shapes = fill_project_shape_objects_flowing(
                    prs,
                    slide,
                    project_shapes,
                    chunks[cursor:],
                    style,
                    allow_oversize_first=not allow_same_slide_full_width,
                    allow_project_splitting=False,
                )
            elif split_layout.masonry_projects:
                used, project_shapes = fill_project_shape_objects_masonry(
                    prs,
                    slide,
                    project_shapes,
                    chunks[cursor:],
                    style,
                    balance_columns=True,
                )
            else:
                used = fill_project_shape_objects_compact(prs, slide, project_shapes, chunks[cursor:], style)
            cursor += used

            # Align the tops of the section and project shapes so the narrow
            # columns start on the same horizontal line.
            common_top = section_shape.top
            for shape in project_shapes:
                common_top = min(common_top, shape.top)
            if section_shape.top != common_top:
                delta = section_shape.top - common_top
                set_shape_geometry(
                    section_shape,
                    section_shape.left,
                    common_top,
                    section_shape.width,
                    section_shape.height + delta,
                )
            for shape in project_shapes:
                if shape.top != common_top:
                    delta = shape.top - common_top
                    set_shape_geometry(
                        shape,
                        shape.left,
                        common_top,
                        shape.width,
                        shape.height + delta,
                    )

            # Make the divider stop at the actual bottom of the rendered
            # narrow-area content, plus a buffer, instead of extending to the
            # fixed layout bottom. This also leaves clear space before the
            # full-width projects below.
            section_bottom = section_shape.top + estimate_filled_shape_height(section_shape, style)
            project_bottom = max(
                (
                    shape.top + estimate_filled_shape_height(shape, style)
                    for shape in project_shapes
                    if shape_has_text(shape)
                ),
                default=section_bottom,
            )
            narrow_bottom = max(section_bottom, project_bottom) + int(0.05 * EMU_PER_INCH)
            narrow_bottom = min(
                split_layout.bottom,
                max(split_layout.top + 100000, narrow_bottom),
            )

            # Resize the section and project shapes to the new bottom so the
            # divider doesn't sit in the middle of empty or overflowing shapes.
            set_shape_geometry(
                section_shape,
                section_shape.left,
                section_shape.top,
                section_shape.width,
                narrow_bottom - section_shape.top,
            )
            for shape in project_shapes:
                set_shape_geometry(
                    shape,
                    shape.left,
                    shape.top,
                    shape.width,
                    narrow_bottom - shape.top,
                )

            if (
                allow_same_slide_full_width
                and cursor < len(chunks)
                and narrow_bottom + int(1.00 * EMU_PER_INCH) < split_layout.bottom
            ):
                # Reserve room for the "Selected Projects" heading above the
                # full-width project lane so it matches the template layout.
                full_top = narrow_bottom + int(0.60 * EMU_PER_INCH)
                full_columns = max(1, split_layout.project_columns)
                full_gap = split_layout.column_gap * (full_columns - 1)
                full_width = max(
                    100000,
                    int(
                        (
                            split_layout.right
                            - split_layout.section_left
                            - full_gap
                        )
                        / full_columns
                    ),
                )
                full_project_shapes: list[Any] = []
                for column in range(full_columns):
                    left = split_layout.section_left + column * (
                        full_width + split_layout.column_gap
                    )
                    full_project_shapes.append(
                        slide.shapes.add_textbox(
                            left,
                            full_top,
                            full_width,
                            split_layout.bottom - full_top,
                        )
                    )
                full_used, _full_flowing_shapes = fill_project_shape_objects_flowing(
                    prs,
                    slide,
                    full_project_shapes,
                    chunks[cursor:],
                    style,
                    allow_oversize_first=False,
                    allow_project_splitting=False,
                    # A short lower-page region is sensitive to PowerPoint's
                    # final justified-line wrapping. Require a 30% safety band
                    # here; otherwise a project that only fits in the estimate
                    # can render below the slide boundary.
                    height_safety_factor=1.30,
                )
                cursor += full_used
                if full_used and full_project_shapes:
                    has_banner = slide_layout_has_teal_banner(slide)
                    add_project_area_heading(
                        slide,
                        full_project_shapes,
                        style,
                        project_only_full_width=True,
                        with_rule=not has_banner,
                    )
            if split_layout.draw_divider:
                add_vertical_divider(
                    slide,
                    split_layout.divider_x,
                    section_shape.top,
                    narrow_bottom,
                )
        else:
            shapes = find_textbox_shapes(slide, shape_indexes)
            shape_cursor = 0
            while remaining_sections and shape_cursor < len(shapes):
                packed, remaining_sections = pack_sections_into_shape(
                    remaining_sections,
                    shapes[shape_cursor],
                    style,
                    keep_sections_together=style is not PORTRAIT_STYLE,
                )
                add_section_text(shapes[shape_cursor], packed, style) if packed else blank_shape(shapes[shape_cursor])
                shape_cursor += 1

            if remaining_sections:
                for shape in shapes[shape_cursor:]:
                    blank_shape(shape)
            elif shape_cursor < len(shapes):
                if style in (PORTRAIT_STYLE, LANDSCAPE_STYLE):
                    used, _flowing_shapes = fill_project_shape_objects_flowing(
                        prs,
                        slide,
                        shapes[shape_cursor:],
                        chunks[cursor:],
                        style,
                    )
                else:
                    used = fill_project_shape_objects_compact(
                        prs,
                        slide,
                        shapes[shape_cursor:],
                        chunks[cursor:],
                        style,
                    )
                cursor += used
                if not remaining_sections and not had_sections_on_slide:
                    max_shape_width = max(s.width for s in shapes[shape_cursor:])
                    is_full_width = max_shape_width >= int(0.8 * prs.slide_width)
                    has_banner = slide_layout_has_teal_banner(slide)
                    add_project_area_heading(
                        slide,
                        shapes[shape_cursor:],
                        style,
                        with_rule=is_full_width and not has_banner,
                    )
            compact_filled_shape_columns(prs, shapes, style)
            cover_layout_labels(slide)
        slide_indexes.append(slide_index)

        if not remaining_sections and cursor >= len(chunks):
            break
        after_sections = sum(len(items) for _heading, items, _continued in remaining_sections)
        if after_sections == before_sections and cursor == before_cursor:
            break

    if not slide_indexes and keep_blank:
        slide = prs.slides[template_index]
        for shape in find_textbox_shapes(slide, shape_indexes):
            blank_shape(shape)
        return [template_index]
    return slide_indexes


def render_section_continuation_slides(
    prs: Presentation,
    template_index: int,
    sections: list[tuple[str, list[str], bool]],
    shape_indexes: list[int],
    style: TextStyle,
    cv: dict[str, Any],
    tmpdir: Path,
) -> list[int]:
    if not sections:
        return []
    template_slide = prs.slides[template_index]
    template_layout = template_slide.slide_layout
    template_elements = [copy.deepcopy(shape.element) for shape in template_slide.shapes]
    slide_indexes: list[int] = []
    remaining = sections
    first = True
    while remaining:
        if first:
            slide = prs.slides[template_index]
            slide_index = template_index
            first = False
        else:
            slide = duplicate_from_snapshot(prs, template_layout, template_elements)
            slide_index = len(prs.slides) - 1
        if len(slide.shapes) > 0 and getattr(slide.shapes[0], "has_text_frame", False):
            replace_photo(slide.shapes[0], cv, tmpdir)
        first_name, last_name, role = profile_parts(cv)
        if len(slide.shapes) > 1:
            fill_name(
                slide.shapes[1], first_name, last_name,
                size=max(16, style.name_size - 2), color=WHITE, font=style.title_font,
            )
        if len(slide.shapes) > 2:
            fill_plain(slide.shapes[2], role, size=style.role_size, color=WHITE, font=style.body_font)
        if len(slide.shapes) > 3:
            fill_plain(slide.shapes[3], "", size=style.overview_size, font=style.body_font)

        shapes = find_textbox_shapes(slide, shape_indexes)
        packed = pack_sections(
            remaining,
            shapes,
            style,
            keep_sections_together=style is not PORTRAIT_STYLE,
        )
        for shape, chunks_for_shape in zip(shapes, packed):
            add_section_text(shape, chunks_for_shape, style) if chunks_for_shape else blank_shape(shape)
        for shape in shapes[len(packed) :]:
            blank_shape(shape)
        for shape in shapes:
            _packed, remaining = pack_sections_into_shape(
                remaining,
                shape,
                style,
                keep_sections_together=style is not PORTRAIT_STYLE,
            )
            if not remaining:
                break
        slide_indexes.append(slide_index)
    return slide_indexes


def render_portrait_section_slides(
    prs: Presentation,
    template_index: int,
    remainders: dict[str, list[str]],
    cv: dict[str, Any],
    tmpdir: Path,
) -> list[int]:
    if not remainders_have_content(remainders):
        return []
    template_slide = prs.slides[template_index]
    template_layout = template_slide.slide_layout
    template_elements = [copy.deepcopy(shape.element) for shape in template_slide.shapes]
    slide_indexes: list[int] = []
    remaining = {key: list(value) for key, value in remainders.items()}
    first = True
    while remainders_have_content(remaining):
        if first:
            slide = prs.slides[template_index]
            slide_index = template_index
            first = False
        else:
            slide = duplicate_from_snapshot(prs, template_layout, template_elements)
            slide_index = len(prs.slides) - 1

        first_name, last_name, role = profile_parts(cv)
        replace_photo(slide.shapes[0], cv, tmpdir)
        fill_name(
            slide.shapes[1], first_name, last_name,
            size=18, color=WHITE, font=PORTRAIT_STYLE.title_font,
        )
        fill_plain(
            slide.shapes[3], role,
            size=PORTRAIT_STYLE.role_size, color=WHITE, font=PORTRAIT_STYLE.body_font,
        )
        fill_plain(slide.shapes[2], "", size=PORTRAIT_STYLE.overview_size, font=PORTRAIT_STYLE.body_font)

        packed, leftover = pack_items_into_shape(remaining.get("key", []), slide.shapes[4], PORTRAIT_STYLE)
        add_items_text(slide.shapes[4], packed, PORTRAIT_STYLE) if packed else blank_shape(slide.shapes[4])
        remaining["key"] = leftover

        qual_blocks = qualification_blocks(remaining)
        packed_blocks, remaining_blocks = pack_blocks_into_shape(qual_blocks, slide.shapes[5], PORTRAIT_STYLE)
        add_block_items_text(slide.shapes[5], packed_blocks, PORTRAIT_STYLE) if packed_blocks else blank_shape(slide.shapes[5])
        apply_qualification_blocks(remaining, remaining_blocks)

        packed, leftover = pack_items_into_shape(remaining.get("experience", []), slide.shapes[6], PORTRAIT_STYLE)
        add_items_text(slide.shapes[6], packed, PORTRAIT_STYLE) if packed else blank_shape(slide.shapes[6])
        remaining["experience"] = leftover

        for shape_index in (7, 8, 9):
            if shape_index < len(slide.shapes):
                blank_shape(slide.shapes[shape_index])
        slide_indexes.append(slide_index)
    return slide_indexes


def render_portrait(cv: dict[str, Any], output_path: Path) -> None:
    configure_experience_body_start(PORTRAIT_STYLE, cv)
    prs = Presentation(str(PORTRAIT_TEMPLATE))
    suppress_layout_labels(prs)
    suppress_portrait_vertical_layout_dividers(prs)
    continuation_bottom = prs.slide_height - int(0.38 * EMU_PER_INCH)
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        option_a_split = SplitContinuationLayout(
            section_left=int(0.50 * EMU_PER_INCH),
            divider_x=int(2.64 * EMU_PER_INCH),
            project_left=int(2.75 * EMU_PER_INCH),
            top=int(1.20 * EMU_PER_INCH),
            bottom=continuation_bottom,
            right=int(6.96 * EMU_PER_INCH),
            project_columns=1,
            draw_divider=True,
            project_only_full_width=True,
        )
        option_b_split = SplitContinuationLayout(
            section_left=int(0.50 * EMU_PER_INCH),
            divider_x=int(2.64 * EMU_PER_INCH),
            project_left=int(2.75 * EMU_PER_INCH),
            top=int(1.20 * EMU_PER_INCH),
            bottom=continuation_bottom,
            right=int(6.96 * EMU_PER_INCH),
            project_columns=1,
            draw_divider=False,
        )
        option_a_sections, option_a_projects = render_front_option(
            prs,
            prs.slides[0],
            cv,
            tmpdir,
            PORTRAIT_STYLE,
            {"photo": 0, "name": 1, "overview": 2, "role": 3},
            {"key": 4, "qualifications": 5, "experience": 6},
            [7, 8],
            flowing_side_sections=True,
            expand_experience=True,
            precompact_before_sections=True,
            front_divider_x=int(2.64 * EMU_PER_INCH),
            front_divider_bottom=int(10.35 * EMU_PER_INCH),
        )
        option_a_slide_indexes = render_mixed_continuation_slides(
            prs,
            1,
            option_a_sections,
            option_a_projects,
            [0, 1, 2, 3, 4, 5],
            PORTRAIT_STYLE,
            keep_blank=True,
            split_layout=option_a_split,
        )
        option_b_sections, option_b_projects = render_front_option(
            prs,
            prs.slides[2],
            cv,
            tmpdir,
            PORTRAIT_STYLE,
            {"photo": 7, "name": 0, "overview": 1, "role": 2},
            {"key": 3, "qualifications": 4, "experience": 5},
            [8, 6],
            name_color=AESG_TEAL,
            role_color=TEXT_DARK,
            photo_background=WHITE,
            flowing_side_sections=True,
            expand_experience=True,
            precompact_before_sections=True,
        )
        option_b_slide_indexes = render_mixed_continuation_slides(
            prs,
            3,
            option_b_sections,
            option_b_projects,
            [0, 1, 2, 3, 4, 5],
            PORTRAIT_STYLE,
            keep_blank=True,
            split_layout=option_b_split,
        )
        assert_project_title_coverage(
            prs,
            cv,
            [0, *option_a_slide_indexes],
            "portrait format 1",
        )
        assert_project_title_coverage(
            prs,
            cv,
            [2, *option_b_slide_indexes],
            "portrait format 2",
        )
        reorder_slides(prs, [0, *option_a_slide_indexes, 2, *option_b_slide_indexes])
        finalize_presentation(prs, tmpdir)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def render_front_landscape(prs: Presentation, cv: dict[str, Any], tmpdir: Path) -> tuple[dict[str, list[str]], list[ProjectChunk]]:
    slide = prs.slides[0]
    first, last, role = profile_parts(cv)
    replace_photo(slide.shapes[0], cv, tmpdir)
    fill_name(
        slide.shapes[1],
        first,
        last,
        size=LANDSCAPE_STYLE.name_size,
        color=WHITE,
        font=LANDSCAPE_STYLE.title_font,
    )
    fill_plain(
        slide.shapes[2],
        role,
        size=LANDSCAPE_STYLE.role_size,
        color=WHITE,
        font=LANDSCAPE_STYLE.body_font,
    )
    fill_plain(
        slide.shapes[3],
        cv.get("overview") or cv.get("overview_full") or "",
        size=LANDSCAPE_STYLE.overview_size,
        font=LANDSCAPE_STYLE.body_font,
        align=PP_ALIGN.JUSTIFY,
    )

    remaining_sections = standard_remainders(cv)
    packed, remaining = pack_items_into_shape(remaining_sections["experience"], slide.shapes[7], LANDSCAPE_STYLE)
    add_items_text(slide.shapes[7], packed, LANDSCAPE_STYLE) if packed else blank_shape(slide.shapes[7])
    remaining_sections["experience"] = remaining

    packed, remaining = pack_items_into_shape(remaining_sections["key"], slide.shapes[4], LANDSCAPE_STYLE)
    add_items_text(slide.shapes[4], packed, LANDSCAPE_STYLE) if packed else blank_shape(slide.shapes[4])
    remaining_sections["key"] = remaining

    qual_blocks = qualification_blocks(remaining_sections)
    packed_blocks, remaining_blocks = pack_blocks_into_shape(qual_blocks, slide.shapes[5], LANDSCAPE_STYLE)
    add_block_items_text(slide.shapes[5], packed_blocks, LANDSCAPE_STYLE) if packed_blocks else blank_shape(slide.shapes[5])
    apply_qualification_blocks(remaining_sections, remaining_blocks)

    project_shapes = [slide.shapes[6], slide.shapes[8], slide.shapes[9], slide.shapes[10]]
    chunks = project_chunks_for_shapes(selected_projects(cv), project_shapes[0], LANDSCAPE_STYLE)
    used, _ = fill_project_shapes(project_shapes, chunks, LANDSCAPE_STYLE)
    return remaining_sections, chunks[used:]


def render_landscape_section_slide(
    prs: Presentation,
    template_index: int,
    sections: dict[str, list[str]],
    cv: dict[str, Any],
    tmpdir: Path,
) -> list[int]:
    if not remainders_have_content(sections):
        return []
    template_slide = prs.slides[template_index]
    template_layout = template_slide.slide_layout
    template_elements = [copy.deepcopy(shape.element) for shape in template_slide.shapes]
    slide_indexes: list[int] = []
    remaining = {key: list(value) for key, value in sections.items()}
    first = True
    while remainders_have_content(remaining):
        if first:
            slide = prs.slides[template_index]
            slide_index = template_index
            first = False
        else:
            slide = duplicate_from_snapshot(prs, template_layout, template_elements)
            slide_index = len(prs.slides) - 1
        replace_photo(slide.shapes[0], cv, tmpdir)
        first_name, last_name, role = profile_parts(cv)
        fill_name(
            slide.shapes[1],
            first_name,
            last_name,
            size=18,
            color=WHITE,
            font=LANDSCAPE_STYLE.title_font,
        )
        fill_plain(
            slide.shapes[2],
            role,
            size=LANDSCAPE_STYLE.role_size,
            color=WHITE,
            font=LANDSCAPE_STYLE.body_font,
        )
        fill_plain(
            slide.shapes[3],
            "",
            size=LANDSCAPE_STYLE.overview_size,
            font=LANDSCAPE_STYLE.body_font,
        )

        packed, leftover = pack_items_into_shape(remaining.get("key", []), slide.shapes[4], LANDSCAPE_STYLE)
        add_items_text(slide.shapes[4], packed, LANDSCAPE_STYLE) if packed else blank_shape(slide.shapes[4])
        remaining["key"] = leftover

        packed, leftover = pack_items_into_shape(remaining.get("experience", []), slide.shapes[5], LANDSCAPE_STYLE)
        add_items_text(slide.shapes[5], packed, LANDSCAPE_STYLE) if packed else blank_shape(slide.shapes[5])
        remaining["experience"] = leftover

        qual_blocks = qualification_blocks(remaining)
        packed_blocks, remaining_blocks = pack_blocks_into_shape(qual_blocks, slide.shapes[6], LANDSCAPE_STYLE)
        add_block_items_text(slide.shapes[6], packed_blocks, LANDSCAPE_STYLE) if packed_blocks else blank_shape(slide.shapes[6])
        apply_qualification_blocks(remaining, remaining_blocks)

        blank_shape(slide.shapes[7])
        slide_indexes.append(slide_index)
    return slide_indexes


def render_landscape(cv: dict[str, Any], output_path: Path) -> None:
    configure_experience_body_start(LANDSCAPE_STYLE, cv)
    prs = Presentation(str(LANDSCAPE_TEMPLATE))
    suppress_layout_labels(prs)
    suppress_portrait_vertical_layout_dividers(prs)
    continuation_bottom = prs.slide_height - int(0.34 * EMU_PER_INCH)
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        option_a_split = SplitContinuationLayout(
            section_left=int(0.50 * EMU_PER_INCH),
            divider_x=int(3.73 * EMU_PER_INCH),
            project_left=int(3.83 * EMU_PER_INCH),
            top=int(1.14 * EMU_PER_INCH),
            bottom=continuation_bottom,
            right=int(10.29 * EMU_PER_INCH),
            project_columns=2,
            column_gap=int(0.20 * EMU_PER_INCH),
            draw_divider=True,
            project_only_full_width=True,
            masonry_projects=True,
            section_density_factor=124.0,
        )
        option_b_split = SplitContinuationLayout(
            section_left=int(0.50 * EMU_PER_INCH),
            divider_x=int(2.64 * EMU_PER_INCH),
            project_left=int(2.75 * EMU_PER_INCH),
            top=int(1.14 * EMU_PER_INCH),
            bottom=continuation_bottom,
            right=int(10.29 * EMU_PER_INCH),
            project_columns=1,
            column_gap=int(0.20 * EMU_PER_INCH),
            draw_divider=False,
            project_only_full_width=True,
            masonry_projects=True,
            section_density_factor=124.0,
        )
        option_a_sections, option_a_projects = render_front_option(
            prs,
            prs.slides[0],
            cv,
            tmpdir,
            LANDSCAPE_STYLE,
            {"photo": 0, "name": 1, "role": 2, "overview": 3},
            {"key": 4, "qualifications": 5, "experience": 7},
            [6, 8, 9, 10],
            flowing_all_sections=True,
            precompact_before_sections=True,
            compact_projects_with_sections=False,
            section_density_factor=124.0,
            project_masonry=True,
            balance_section_overflow=False,
            align_project_heading_with_sections=True,
            front_divider_x=int(3.73 * EMU_PER_INCH),
            front_divider_bottom=int(7.16 * EMU_PER_INCH),
            section_bottom_buffer=int(0.15 * EMU_PER_INCH),
        )
        option_a_slide_indexes = render_mixed_continuation_slides(
            prs,
            1,
            option_a_sections,
            option_a_projects,
            [0, 3, 1, 4, 2, 5],
            LANDSCAPE_STYLE,
            keep_blank=True,
            section_shape_indexes=[0, 1, 2],
            project_shape_indexes=[3, 4, 5],
            split_layout=option_a_split,
        )
        option_b_sections, option_b_projects = render_front_option(
            prs,
            prs.slides[2],
            cv,
            tmpdir,
            LANDSCAPE_STYLE,
            {"photo": 0, "name": 1, "role": 2, "overview": 3},
            {"key": 4, "experience": 5, "qualifications": 6},
            [7],
            name_color=AESG_TEAL,
            role_color=TEXT_DARK,
            photo_background=WHITE,
            flowing_side_sections=True,
            expand_experience=True,
            precompact_before_sections=True,
            content_respects_photo=True,
            section_density_factor=124.0,
            section_bottom_buffer=int(0.15 * EMU_PER_INCH),
        )
        option_b_slide_indexes = render_mixed_continuation_slides(
            prs,
            3,
            option_b_sections,
            option_b_projects,
            [0, 1, 2],
            LANDSCAPE_STYLE,
            keep_blank=True,
            split_layout=option_b_split,
        )
        assert_project_title_coverage(
            prs,
            cv,
            [0, *option_a_slide_indexes],
            "landscape format 1",
        )
        assert_project_title_coverage(
            prs,
            cv,
            [2, *option_b_slide_indexes],
            "landscape format 2",
        )
        reorder_slides(prs, [0, *option_a_slide_indexes, 2, *option_b_slide_indexes])
        finalize_presentation(prs, tmpdir)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def assert_project_title_coverage(
    prs: Presentation,
    cv: dict[str, Any],
    slide_indexes: list[int],
    format_label: str,
) -> None:
    """Fail generation if either PPTX format silently drops a project."""
    rendered_paragraphs = [
        clean_text(paragraph_text(paragraph))
        for slide_index in slide_indexes
        if 0 <= slide_index < len(prs.slides)
        for shape in prs.slides[slide_index].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if clean_text(paragraph_text(paragraph))
    ]
    expected_titles = Counter(
        clean_text(project.get("name") or "")
        for project in selected_projects(cv)
        if clean_text(project.get("name") or "")
    )
    missing: list[str] = []
    for title, expected_count in expected_titles.items():
        rendered_count = sum(
            paragraph == title
            or paragraph.startswith(f"{title} (continued ")
            for paragraph in rendered_paragraphs
        )
        missing.extend([title] * max(0, expected_count - rendered_count))
    if missing:
        raise ValueError(
            f"{format_label} omitted {len(missing)} selected project(s): "
            + "; ".join(missing)
        )


def validate_deck(path: Path) -> list[str]:
    warnings: list[str] = []
    prs = Presentation(str(path))
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = clean_text(shape.text)
            if "{{" in text or "}}" in text:
                warnings.append(f"{path.name}: unresolved tag on slide {slide_index} shape {shape_index}")
            if not text:
                continue
            size = 7.0
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        size = run.font.size.pt
                        break
                if size:
                    break
            paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph_text(paragraph)]
            has_bullets = any(
                paragraph._p.pPr is not None
                and any(child.tag.rsplit("}", 1)[-1] in {"buChar", "buAutoNum", "buBlip"} for child in paragraph._p.pPr)
                for paragraph in paragraphs
            )
            has_dark_body = any(
                str(run.font.color.rgb) == TEXT_DARK
                for paragraph in paragraphs[1:]
                for run in paragraph.runs
                if run.font.color is not None
            )
            if len(paragraphs) > 1 and not has_bullets and has_dark_body:
                if prs.slide_height > prs.slide_width:
                    estimated = estimate_project_filled_shape_height(shape, PORTRAIT_STYLE) / 12700
                else:
                    estimated = estimate_compact_project_filled_height(shape, LANDSCAPE_STYLE) / 12700
            else:
                # Long justified overview paragraphs render at a substantially
                # higher character density than bullets or narrow body cards.
                # Applying the generic bullet estimate reports false overflow
                # even when PowerPoint leaves visible space below the text.
                factor = 140.0 if len(paragraphs) == 1 and len(text) > 350 else 98.0
                leading = 1.16
                if len(paragraphs) == 1:
                    paragraph_leading = paragraphs[0].line_spacing
                    if isinstance(paragraph_leading, (int, float)):
                        leading = float(paragraph_leading)
                estimated = paragraph_height(
                    text,
                    shape.width,
                    size,
                    factor=factor,
                    leading=leading,
                )
            capacity = available_height_pt(shape, padding_pt=0)
            if estimated > capacity * 1.18:
                warnings.append(f"{path.name}: possible overflow slide {slide_index} shape {shape_index}")
    return warnings


def ensure_templates() -> None:
    if PORTRAIT_TEMPLATE.exists() and LANDSCAPE_TEMPLATE.exists():
        return
    import importlib.util

    builder_path = ROOT / "scripts" / "v3" / "build_pptx_templates_v3.py"
    spec = importlib.util.spec_from_file_location("build_pptx_templates_v3", builder_path)
    if not spec or not spec.loader:
        raise RuntimeError("Could not load PPTX template builder")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    module.main()


def cv_for_rendering(cv: dict[str, Any]) -> dict[str, Any]:
    """Keep legacy Education fields empty when rendering older payloads."""
    rendered = copy.deepcopy(cv)
    rendered["education"] = []
    rendered["education_portrait"] = []
    rendered["education_landscape"] = []
    return rendered


def main() -> int:
    parser = argparse.ArgumentParser(description="Render V3 PPTX CV outputs.")
    parser.add_argument("--data", type=Path, default=DATA_PATH)
    parser.add_argument(
        "--output-root",
        type=Path,
        default=ROOT / "output" / "generated",
        help="Root folder containing mirrored source folders and PPTX artifact folders.",
    )
    parser.add_argument("--limit", type=int, default=None)
    parser.add_argument("--orientation", choices=("portrait", "landscape", "both"), default="both")
    args = parser.parse_args()

    ensure_templates()
    cvs = load_cvs(args.data)
    if args.limit:
        cvs = cvs[: args.limit]

    warnings: list[str] = []
    claimed_paths: set[Path] = set()
    for cv in cvs:
        rendered_cv = cv_for_rendering(cv)
        if args.orientation in {"portrait", "both"}:
            out = artifact_output_path(
                args.output_root,
                cv,
                "pptx_portrait",
                ".pptx",
            )
            if out in claimed_paths:
                raise ValueError(
                    f"Duplicate source PDF path would overwrite an output: {out}"
                )
            claimed_paths.add(out)
            render_portrait(rendered_cv, out)
            print(out)
            warnings.extend(validate_deck(out))
        if args.orientation in {"landscape", "both"}:
            out = artifact_output_path(
                args.output_root,
                cv,
                "pptx_landscape",
                ".pptx",
            )
            if out in claimed_paths:
                raise ValueError(
                    f"Duplicate source PDF path would overwrite an output: {out}"
                )
            claimed_paths.add(out)
            render_landscape(rendered_cv, out)
            print(out)
            warnings.extend(validate_deck(out))

    if warnings:
        print("\nValidation warnings:")
        for warning in warnings:
            print(f"- {warning}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
