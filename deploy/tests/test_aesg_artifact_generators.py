from __future__ import annotations

import importlib.util
import tempfile
import unittest
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches


REPOSITORY = Path(__file__).resolve().parents[2]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"unable to load {path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


validator = load_module(
    "aesg_validate_artifact",
    REPOSITORY / "deploy/skills/aesg-branding/scripts/validate_artifact.py",
)
docx_generator = load_module(
    "aesg_create_docx",
    REPOSITORY / "deploy/skills/docx/scripts/create_aesg_docx.py",
)
docx_repair = load_module(
    "aesg_repair_docx",
    REPOSITORY / "deploy/skills/docx/scripts/repair_aesg_docx.py",
)


class ArtifactPlaceholderValidationTests(unittest.TestCase):
    def test_docx_placeholder_split_across_runs_is_detected(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "split-placeholder.docx"
            document = Document()
            paragraph = document.add_paragraph()
            paragraph.add_run("Lorem")
            paragraph.add_run(" ipsum").bold = True
            document.save(path)

            text = validator.docx_story_text(path)

        self.assertIn("Lorem ipsum", text)
        self.assertIsNotNone(validator.PLACEHOLDER_RE.search(text))

    def test_pptx_placeholder_split_across_runs_is_detected(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "split-placeholder.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_frame = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(4), Inches(1)
            ).text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            paragraph.add_run().text = "Lorem"
            paragraph.add_run().text = " ipsum"
            presentation.save(path)
            slide_order = {
                str(item.part.partname).lstrip("/"): index
                for index, item in enumerate(presentation.slides, start=1)
            }

            text, _, _ = validator.inspect_pptx_slides(path, slide_order)

        self.assertIn("Lorem ipsum", text)
        self.assertIsNotNone(validator.PLACEHOLDER_RE.search(text))


class LetterheadGeometryTests(unittest.TestCase):
    def test_tables_fit_retained_letterhead_geometry_and_keep_type_size(self) -> None:
        template = (
            REPOSITORY
            / "deploy/skills/aesg-branding/assets/templates/AESG_Letterhead_Dubai.docx"
        )
        document = Document(template)
        docx_generator.create_letter(
            document,
            {
                "title": "Leave request",
                "recipient": ["People Team", "AESG"],
                "subject": "Annual leave request",
                "sections": [
                    {
                        "paragraphs": ["Please approve the requested leave dates."],
                        "callout": "Requested dates are subject to manager approval.",
                        "table": {
                            "headers": ["From", "To", "Working days"],
                            "widths": [2, 2, 1],
                            "rows": [["20 August 2026", "24 August 2026", "3"]],
                        },
                    }
                ],
                "signatory": "AESG employee",
            },
        )
        docx_generator.normalise_new_text(document)

        expected_content_width = 9_072
        self.assertEqual(docx_generator.usable_width_dxa(document), expected_content_width)
        self.assertEqual(len(document.tables), 2)
        for table in document.tables:
            properties = table._tbl.tblPr
            width = int(properties.find(qn("w:tblW")).get(qn("w:w")))
            indent = int(properties.find(qn("w:tblInd")).get(qn("w:w")))
            grid_width = sum(
                int(column.get(qn("w:w"))) for column in table._tbl.tblGrid
            )
            self.assertEqual(width, expected_content_width - docx_generator.CELL_MARGIN_DXA)
            self.assertEqual(grid_width, width)
            self.assertEqual(width + indent, expected_content_width)
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if not run.text:
                                continue
                            self.assertEqual(run.font.name, docx_generator.FONT)
                            self.assertAlmostEqual(
                                run.font.size.pt,
                                docx_generator.LETTER_BODY_PT,
                            )


class ReportStructureTests(unittest.TestCase):
    def setUp(self) -> None:
        self.template = (
            REPOSITORY
            / "deploy/skills/aesg-branding/assets/templates/AESG_General_Report_Template.docx"
        )

    def test_ordered_blocks_preserve_sequence_and_real_list_styles(self) -> None:
        document = Document(self.template)
        docx_generator.create_report(
            document,
            {
                "title": "Structure smoke test",
                "approvalPage": False,
                "sections": [
                    {
                        "heading": "Actions",
                        "blocks": [
                            {"type": "paragraph", "text": "Opening prose."},
                            {"type": "label", "text": "Recommendations:"},
                            {"type": "bullets", "items": ["First action", "Second action"]},
                            {"type": "paragraph", "text": "Closing prose."},
                        ],
                    }
                ],
            },
        )

        by_text = {paragraph.text: paragraph for paragraph in document.paragraphs}
        ordered = [
            next(
                index
                for index, paragraph in enumerate(document.paragraphs)
                if paragraph.text == text
            )
            for text in (
                "Opening prose.",
                "Recommendations:",
                "First action",
                "Second action",
                "Closing prose.",
            )
        ]
        self.assertEqual(ordered, sorted(ordered))
        self.assertEqual(by_text["First action"].style.name, "AESG Bullet")
        self.assertEqual(by_text["Second action"].style.name, "AESG Bullet")
        self.assertTrue(by_text["Recommendations:"].runs[0].bold)

    def test_monitoring_table_block_uses_a_landscape_section(self) -> None:
        document = Document(self.template)
        docx_generator.create_report(
            document,
            {
                "title": "Monitoring layout test",
                "approvalPage": False,
                "sections": [
                    {
                        "heading": "Monitoring summary",
                        "blocks": [
                            {
                                "type": "table",
                                "table": {
                                    "pattern": "monitoring",
                                    "headers": ["Location", "Reading", "Status"],
                                    "rows": [["North boundary", "42", "Compliant"]],
                                },
                            }
                        ],
                    }
                ],
            },
        )

        self.assertTrue(
            any(section.page_width > section.page_height for section in document.sections)
        )

    def test_repair_disables_template_numbering_for_manual_heading(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "manual-heading.docx"
            document = Document(self.template)
            docx_generator.create_report(
                document,
                {
                    "title": "Numbering smoke test",
                    "approvalPage": False,
                    "sections": [
                        {"heading": "2.1 Background", "level": 2},
                        {
                            "table": {
                                "headers": ["Requirement", "Gap", "RAG"],
                                "rows": [["Access control", "Upgrade required", "RED"]],
                            }
                        },
                    ],
                },
            )
            document.save(path)
            docx_repair.repair_report(path)
            repaired = Document(path)
            heading = next(
                paragraph
                for paragraph in repaired.paragraphs
                if paragraph.text == "2.1 Background"
            )
            numbering = heading._p.pPr.find(qn("w:numPr"))
            self.assertEqual(numbering.find(qn("w:numId")).get(qn("w:val")), "0")
            compliance = next(
                table
                for table in repaired.tables
                if table.rows and table.cell(0, 0).text == "Requirement"
            )
            self.assertEqual(compliance.cell(1, 2).text, "Red")

    def test_validator_rejects_flat_lists_and_manual_heading_numbering(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "flat-report.docx"
            document = Document(self.template)
            docx_generator.create_report(
                document,
                {
                    "title": "Validation smoke test",
                    "approvalPage": False,
                    "sections": [
                        {
                            "heading": "2.1 Background",
                            "level": 2,
                            "paragraphs": ["Examples:", "One", "Two", "Three"],
                        }
                    ],
                },
            )
            document.save(path)
            docx_generator.normalise_report_output_parts(path)
            errors: list[str] = []
            evidence: dict = {}
            validator.validate_docx(path, errors, evidence)

        self.assertTrue(any("inherit automatic numbering" in error for error in errors))
        self.assertTrue(any("flattened as body paragraphs" in error for error in errors))


if __name__ == "__main__":
    unittest.main()
